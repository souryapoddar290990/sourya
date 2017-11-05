# tokenizer = None
# tagger = None

# def init_nltk():
#     global tokenizer
#     global tagger
#     tokenizer = nltk.tokenize.RegexpTokenizer(r'\w+|[^\w\s]+')
#     tagger = nltk.UnigramTagger(nltk.corpus.brown.tagged_sents())

# def tag(text):
#     global tokenizer
#     global tagger
#     if not tokenizer:
#         init_nltk()
#     tokenized = tokenizer.tokenize(text)
#     tagged = tagger.tag(tokenized)
#     tagged.sort(lambda x,y:cmp(x[1],y[1]))
#     return tagged

# def main():
#     text = """Mr Blobby is a fictional character who featured on Noel Edmonds' Saturday night entertainment show Noel's House Party,
#     which was often a ratings winner in the 1990s. Mr Blobby also appeared on the Jamie Rose show of 1997. He was designed as an
#     outrageously over the top parody of a one-dimensional, mute novelty character, which ironically made him distinctive, absurd and popular.
#     He was a large pink humanoid, covered with yellow spots, sporting a permanent toothy grin and jiggling eyes. He communicated by saying
#     the word "blobby" in an electronically-altered voice, expressing his moods through tone of voice and repetition.

#     There was a Mrs. Blobby, seen briefly in the video, and sold as a doll.

#     However Mr Blobby actually started out as part of the 'Gotcha' feature during the show's second series (originally called 'Gotcha
#     Oscars' until the threat of legal action from the Academy of Motion Picture Arts and Sciences[citation needed]), in which celebrities
#     were caught out in a Candid Camera style prank. Celebrities such as dancer Wayne Sleep and rugby union player Will Carling would be
#     enticed to take part in a fictitious children's programme based around their profession. Mr Blobby would clumsily take part in the activity,
#     knocking over the set, causing mayhem and saying "blobby blobby blobby", until finally when the prank was revealed, the Blobby
#     costume would be opened - revealing Noel inside. This was all the more surprising for the "victim" as during rehearsals Blobby would be
#     played by an actor wearing only the arms and legs of the costume and speaking in a normal manner.[citation needed]"""
#     tagged = tag(text) 
#     l = list(set(tagged))
#     l.sort(lambda x,y:cmp(x[1],y[1]))
#     # pprint.pprint(l)

# main()

# import speech_recognition as sr
# r = sr.Recognizer()
# with sr.AudioFile("woman1_wb.wav") as source:
#     audio = r.record(source)
# try:
#     s = r.recognize_google(audio)
#     print("Text: "+s)
# except Exception as e:
#     print("Exception: "+str(e))

from scipy.cluster import hierarchy
import editdistance
# def clustering(inputs,inter_cluster_distance):
#     dist, clust_details = [], {}
#     for i in range(len(inputs)-1):
#         for j in range(i+1,len(inputs)): dist.append(int(editdistance.eval(inputs[i], inputs[j])))
#     clusters = hierarchy.linkage(dist, method='weighted', metric='euclidean')
#     T = hierarchy.fcluster(clusters, inter_cluster_distance, criterion='distance')
#     for item in range(max(T)): clust_details[item+1] = []
#     for index,item in enumerate(inputs): clust_details[T[index]].append(item)
#     return clust_details

# inputs = ['banana','bahama','xyz']
# inputs = ['arya','arja','sourya','soury','poddar','soarpa']
# inputs = ['arya','arja','aryan','ayan','rudrashis','indrashis','sourya','surja','sourashis']
# inter_cluster_distance = 3
# print clustering(inputs,inter_cluster_distance)

import pandas as pd
from pandas import Series, DataFrame
import numpy as np
import seaborn as sns
import matplotlib.pyplot as plt
from collections import Counter
import pprint,time,requests,urllib2,json,os,sys,csv,nltk
from nltk import word_tokenize

def cricket_score():
    url = 'http://www.cricapi.com/api/ballByBall?unique_id=1000881&apikey=gWkkyahOxVPSQxn80fBUiz2Nq6P2'
    response = urllib2.urlopen(url)
    result = response.read()
    result = json.loads(result)
    team = result['team']
    data = result['data']
    player_ids = {}
    team_ids = {}
    team1 = team[0]
    team2 = team[1]
    print team1['team_short_name'],'vs',team2['team_short_name']
    team_ids[str(team1['team_id'])] = team1['team_filename']
    team_ids[str(team2['team_id'])] = team2['team_filename']
    print '\n'
    print team1['team_short_name'],'XI'
    for item in team1['player']:
        text = item['known_as']
        if int(item['keeper']) == 1: text += ' ('+'W)'
        if int(item['captain']) == 1: text += ' ('+'C)'
        player_ids[str(item['player_id'])] = item['card_long']
        print text
    print '\n'
    print team2['team_short_name'],'XI'
    for item in team2['player']:
        text = item['known_as']
        if int(item['keeper']) == 1: text += ' ('+'W)'
        if int(item['captain']) == 1: text += ' ('+'C)'
        player_ids[str(item['player_id'])] = item['card_long']
        print text
    print '\nSHORT SCORE\n'
    print team_ids[str(data[1]['team_id'])],str(data[1]['runs']),'/',str(data[1]['wickets']),'(',str(data[1]['over_number']),')'
    print '\n'
    bat1 = data[1]['batsman'][0]
    bat2 = data[1]['batsman'][1]
    bowl1 = data[1]['bowler'][0]
    bowl2 = data[1]['bowler'][1]
    print player_ids[str(bat1['old_player_id'])],bat1['runs'],bat1['balls_faced'],bat1['fours'],bat1['sixes']
    print player_ids[str(bat2['old_player_id'])],bat2['runs'],bat2['balls_faced'],bat2['fours'],bat2['sixes']
    print '\n'
    print player_ids[str(bowl1['old_player_id'])],bowl1['overs'],bowl1['maidens'],bowl1['conceded'],bowl1['wickets']
    print player_ids[str(bowl2['old_player_id'])],bowl2['overs'],bowl2['maidens'],bowl2['conceded'],bowl2['wickets']
    print '\n'
# cricket_score()

# from nltk.corpus import wordnet as wn
# from pyxdameraulevenshtein import damerau_levenshtein_distance, normalized_damerau_levenshtein_distance
# import editdistance,ast,pprint

# sub_data = {'gw': 1, 'gv': 0, 'gu': 0, 'gt': 21, 'gs': 13, 'gr': 5, 'gq': 3, 'gp': 1, 'gz': 0, 'gy': 3, 'gx': 0, 'gg': 0, 'gf': 2, 'ge': 9, 'gd': 11, 'gc': 11, 'gb': 1, 'ga': 4, 'go': 2, 'gn': 0, 'gm': 0, 'gl': 3, 'gk': 1, 'gj': 1, 'gi': 0, 'gh': 0, 'tz': 6, 'tx': 0, 'ty': 7, 'tv': 2, 'tw': 19, 'tt': 0, 'tu': 0, 'tr': 11, 'ts': 37, 'tp': 6, 'tq': 0, 'tn': 5, 'to': 5, 'tl': 14, 'tm': 9, 'tj': 1, 'tk': 0, 'th': 5, 'ti': 0, 'tf': 5, 'tg': 19, 'td': 42, 'te': 7, 'tb': 4, 'tc': 9, 'ta': 3, 'vu': 0, 'zl': 7, 'zm': 5, 'zn': 0, 'zo': 0, 'zh': 0, 'zi': 0, 'zj': 0, 'zk': 0, 'zd': 7, 'ze': 0, 'zf': 0, 'zg': 0, 'za': 0, 'zb': 0, 'zc': 0, 'zx': 0, 'zy': 3, 'zz': 0, 'zt': 3, 'zu': 0, 'zv': 0, 'zw': 0, 'zp': 0, 'zq': 0, 'zr': 2, 'zs': 21, 'wl': 0, 'va': 0, 'vc': 7, 'wk': 1, 'vh': 0, 'wj': 0, 'vi': 0, 'vj': 0, 'vk': 0, 'vl': 1, 'vm': 0, 'wi': 0, 'vn': 0, 'vo': 1, 'me': 0, 'md': 8, 'mg': 0, 'mf': 2, 'ma': 1, 'mc': 7, 'mb': 3, 'mm': 0, 'ml': 4, 'mo': 0, 'mn': 180, 'mi': 0, 'mh': 6, 'mk': 4, 'mj': 0, 'mu': 13, 'mt': 15, 'mw': 2, 'mv': 3, 'mq': 0, 'mp': 6, 'ms': 9, 'mr': 0, 'vt': 3, 'my': 3, 'mx': 2, 'mz': 0, 'vv': 0, 'vw': 0, 'vx': 0, 'vz': 0, 'fp': 0, 'fq': 0, 'fr': 6, 'fs': 4, 'ft': 12, 'fu': 0, 'fv': 0, 'fw': 2, 'fx': 0, 'fy': 0, 'fz': 0, 'fa': 0, 'fb': 15, 'fc': 0, 'fd': 3, 'fe': 1, 'ff': 0, 'fg': 5, 'fh': 2, 'fi': 0, 'fj': 0, 'fk': 0, 'fl': 3, 'fm': 4, 'fn': 1, 'fo': 0, 'sz': 1, 'sy': 20, 'sx': 3, 'ss': 0, 'sr': 14, 'sq': 0, 'sp': 7, 'sw': 5, 'sv': 0, 'su': 0, 'st': 15, 'sk': 0, 'sj': 1, 'si': 0, 'sh': 1, 'so': 1, 'sn': 6, 'sm': 0, 'sl': 27, 'sc': 27, 'sb': 8, 'sa': 11, 'sg': 0, 'sf': 4, 'se': 35, 'sd': 33, 'lf': 4, 'lg': 5, 'ld': 4, 'le': 0, 'lb': 10, 'lc': 1, 'la': 2, 'ln': 14, 'lo': 2, 'll': 0, 'lm': 0, 'lj': 0, 'lk': 1, 'lh': 6, 'li': 13, 'lv': 0, 'lw': 0, 'lt': 2, 'lu': 0, 'lr': 11, 'ls': 10, 'lp': 5, 'lq': 0, 'lz': 0, 'lx': 0, 'ly': 0, 'wq': 0, 'yh': 7, 'yk': 0, 'yj': 0, 'ym': 2, 'yl': 0, 'yo': 6, 'yn': 0, 'ya': 0, 'yc': 2, 'yb': 0, 'ye': 15, 'yd': 0, 'yg': 1, 'yf': 0, 'yy': 0, 'yx': 1, 'yz': 0, 'yq': 0, 'yp': 1, 'ys': 36, 'yr': 7, 'yu': 5, 'yt': 8, 'yw': 0, 'yv': 0, 'em': 0, 'el': 3, 'eo': 93, 'en': 5, 'ei': 89, 'eh': 0, 'ek': 0, 'ej': 0, 'ee': 0, 'ed': 11, 'eg': 2, 'ef': 2, 'ea': 388, 'ec': 3, 'eb': 0, 'ey': 18, 'ex': 0, 'ez': 0, 'eu': 15, 'et': 6, 'ew': 1, 'ev': 0, 'eq': 0, 'ep': 0, 'es': 12, 'er': 14, 'rt': 22, 'ru': 4, 'rv': 0, 'rw': 0, 'rp': 14, 'rq': 0, 'rr': 0, 'rs': 12, 'rx': 1, 'ry': 0, 'rz': 0, 'rd': 30, 're': 12, 'rf': 2, 'rg': 2, 'ra': 0, 'rb': 14, 'rc': 0, 'rl': 8, 'rm': 4, 'rn': 20, 'ro': 1, 'rh': 8, 'ri': 2, 'rj': 0, 'rk': 5, 'xj': 0, 'xk': 0, 'xh': 0, 'xi': 0, 'xn': 0, 'xo': 0, 'xl': 0, 'xm': 0, 'xb': 0, 'xc': 0, 'xa': 0, 'xf': 0, 'xg': 0, 'xd': 2, 'xe': 0, 'xz': 0, 'xx': 0, 'xy': 0, 'xr': 0, 'xs': 9, 'xp': 0, 'xq': 0, 'xv': 0, 'xw': 0, 'xt': 0, 'xu': 0, 'wy': 0, 'wx': 0, 'kc': 8, 'kb': 2, 'ka': 1, 'kg': 2, 'kf': 1, 'ke': 1, 'kd': 4, 'kk': 0, 'kj': 0, 'ki': 0, 'kh': 5, 'ko': 2, 'kn': 0, 'km': 5, 'kl': 0, 'ks': 6, 'kr': 0, 'kq': 0, 'kp': 0, 'kw': 4, 'kv': 0, 'ku': 0, 'kt': 0, 'kz': 3, 'ky': 0, 'kx': 0, 'dn': 3, 'do': 0, 'dl': 3, 'dm': 7, 'dj': 0, 'dk': 2, 'dh': 5, 'di': 0, 'df': 0, 'dg': 5, 'dd': 0, 'de': 12, 'db': 10, 'dc': 13, 'da': 1, 'dz': 0, 'dx': 0, 'dy': 2, 'dv': 0, 'dw': 4, 'dt': 22, 'du': 0, 'dr': 43, 'ds': 30, 'dp': 1, 'dq': 0, 'qq': 0, 'qp': 0, 'qs': 0, 'qr': 0, 'qu': 0, 'qt': 0, 'qw': 0, 'qv': 0, 'qy': 0, 'qx': 0, 'qz': 0, 'qa': 0, 'qc': 1, 'qb': 0, 'qe': 0, 'qd': 0, 'qg': 27, 'qf': 0, 'qi': 0, 'qh': 0, 'qk': 0, 'qj': 0, 'qm': 0, 'ql': 0, 'qo': 0, 'qn': 0, 'wc': 1, 'wb': 2, 'wa': 2, 'wo': 0, 'wn': 0, 'wm': 0, 'wg': 0, 'wf': 0, 'we': 1, 'wd': 0, 'jx': 0, 'jy': 0, 'jz': 0, 'jt': 0, 'ju': 0, 'jv': 0, 'jw': 0, 'jp': 0, 'jq': 0, 'jr': 0, 'js': 5, 'jl': 2, 'jm': 1, 'jn': 0, 'jo': 0, 'jh': 0, 'ji': 0, 'jj': 0, 'jk': 0, 'jd': 9, 'je': 0, 'jf': 0, 'jg': 1, 'ja': 0, 'jb': 1, 'jc': 1, 'ww': 0, 'wv': 0, 'wu': 1, 'wt': 3, 'ws': 3, 'wr': 6, 'ck': 1, 'cj': 0, 'ci': 0, 'ch': 0, 'co': 1, 'cn': 9, 'cm': 7, 'cl': 0, 'cc': 0, 'cb': 5, 'ca': 6, 'wp': 7, 'cg': 5, 'cf': 9, 'ce': 0, 'cd': 16, 'cz': 0, 'cy': 1, 'cx': 1, 'cs': 39, 'cr': 5, 'cq': 2, 'cp': 10, 'cw': 7, 'cv': 3, 'cu': 1, 'ct': 40, 'pr': 1, 'ps': 3, 'pp': 0, 'pq': 0, 'pv': 4, 'pw': 1, 'pt': 6, 'pu': 0, 'pz': 0, 'px': 0, 'py': 0, 'wz': 0, 'pb': 11, 'pc': 1, 'pa': 0, 'pf': 6, 'pg': 5, 'pd': 2, 'pe': 0, 'pj': 9, 'pk': 0, 'ph': 0, 'pi': 2, 'pn': 6, 'po': 15, 'pl': 2, 'pm': 7, 'iy': 15, 'ix': 1, 'vb': 0, 'iz': 0, 'vd': 0, 've': 0, 'vf': 3, 'vg': 0, 'iq': 0, 'ip': 0, 'is': 2, 'ir': 0, 'iu': 47, 'it': 1, 'iw': 2, 'iv': 0, 'ii': 0, 'ih': 0, 'ik': 0, 'ij': 0, 'im': 0, 'il': 6, 'io': 49, 'in': 0, 'ia': 103, 'vy': 0, 'ic': 0, 'ib': 0, 'ie': 146, 'id': 0, 'ig': 1, 'if': 0, 'wh': 2, 'yi': 15, 'vr': 0, 'vs': 8, 'bd': 9, 'be': 2, 'bf': 2, 'bg': 3, 'ba': 0, 'bb': 0, 'bc': 9, 'bl': 5, 'bm': 11, 'bn': 5, 'bo': 0, 'bh': 1, 'bi': 0, 'bj': 0, 'bk': 0, 'bt': 1, 'bu': 0, 'bv': 0, 'bw': 8, 'bp': 10, 'bq': 0, 'br': 0, 'bs': 2, 'bx': 0, 'by': 0, 'bz': 0, 'oo': 0, 'on': 0, 'om': 0, 'ol': 0, 'ok': 2, 'oj': 0, 'oi': 25, 'oh': 0, 'og': 0, 'of': 0, 'oe': 116, 'od': 3, 'oc': 1, 'ob': 1, 'oa': 91, 'oz': 0, 'oy': 18, 'ox': 0, 'ow': 0, 'ov': 0, 'ou': 39, 'ot': 14, 'os': 4, 'or': 2, 'oq': 0, 'op': 14, 'hz': 0, 'hx': 0, 'hy': 0, 'hr': 3, 'hs': 1, 'hp': 3, 'hq': 0, 'hv': 0, 'hw': 2, 'ht': 11, 'hu': 0, 'hj': 0, 'hk': 2, 'hh': 0, 'hi': 0, 'hn': 14, 'ho': 2, 'hl': 0, 'hm': 12, 'hb': 8, 'hc': 0, 'ha': 1, 'hf': 0, 'hg': 0, 'hd': 3, 'he': 0, 'uy': 8, 'ux': 0, 'uz': 0, 'uu': 0, 'ut': 0, 'uw': 2, 'uv': 0, 'uq': 0, 'up': 0, 'us': 0, 'ur': 4, 'um': 0, 'ul': 0, 'uo': 43, 'un': 2, 'ui': 64, 'uh': 0, 'uk': 0, 'uj': 0, 'ue': 44, 'ud': 0, 'ug': 0, 'uf': 0, 'ua': 20, 'uc': 0, 'ub': 0, 'aa': 0, 'ac': 7, 'ab': 0, 'ae': 342, 'ad': 1, 'ag': 0, 'af': 0, 'ai': 118, 'ah': 2, 'ak': 1, 'aj': 0, 'am': 0, 'al': 0, 'ao': 76, 'an': 3, 'aq': 0, 'ap': 0, 'as': 35, 'ar': 1, 'au': 9, 'at': 9, 'aw': 1, 'av': 0, 'ay': 5, 'ax': 0, 'az': 0, 'nh': 19, 'ni': 1, 'nj': 0, 'nk': 4, 'nl': 35, 'nm': 78, 'nn': 0, 'no': 0, 'na': 2, 'nb': 7, 'nc': 6, 'nd': 5, 'ne': 3, 'nf': 0, 'ng': 1, 'nx': 2, 'ny': 0, 'nz': 2, 'np': 7, 'nq': 0, 'nr': 28, 'ns': 5, 'nt': 7, 'nu': 0, 'nv': 0, 'nw': 1, 'vp': 0, 'vq': 0}
# int_data = {'gw': 0, 'gv': 0, 'gu': 3, 'gt': 0, 'gs': 0, 'gr': 3, 'gq': 0, 'gp': 0, 'gz': 0, 'gy': 0, 'gx': 0, 'gg': 0, 'gf': 0, 'ge': 2, 'gd': 0, 'gc': 0, 'gb': 0, 'ga': 4, 'go': 0, 'gn': 15, 'gm': 0, 'gl': 1, 'gk': 0, 'gj': 0, 'gi': 0, 'gh': 0, 'tz': 0, 'tx': 0, 'ty': 0, 'tv': 0, 'tw': 2, 'tt': 0, 'tu': 11, 'tr': 5, 'ts': 0, 'tp': 0, 'tq': 0, 'tn': 0, 'to': 3, 'tl': 4, 'tm': 0, 'tj': 0, 'tk': 0, 'th': 21, 'ti': 49, 'tf': 0, 'tg': 0, 'td': 0, 'te': 4, 'tb': 0, 'tc': 3, 'ta': 4, 'vu': 0, 'zl': 0, 'zm': 0, 'zn': 0, 'zo': 0, 'zh': 0, 'zi': 0, 'zj': 0, 'zk': 0, 'zd': 0, 'ze': 0, 'zf': 0, 'zg': 0, 'za': 0, 'zb': 0, 'zc': 0, 'zx': 0, 'zy': 0, 'zz': 0, 'zt': 0, 'zu': 0, 'zv': 0, 'zw': 0, 'zp': 0, 'zq': 0, 'zr': 0, 'zs': 0, 'wl': 0, 'va': 0, 'vc': 0, 'wk': 0, 'vh': 0, 'wj': 0, 'vi': 4, 'vj': 0, 'vk': 0, 'vl': 0, 'vm': 0, 'wi': 0, 'vn': 0, 'vo': 0, 'me': 20, 'md': 0, 'mg': 0, 'mf': 0, 'ma': 9, 'mc': 0, 'mb': 0, 'mm': 0, 'ml': 0, 'mo': 0, 'mn': 2, 'mi': 0, 'mh': 0, 'mk': 0, 'mj': 0, 'mu': 4, 'mt': 0, 'mw': 0, 'mv': 0, 'mq': 0, 'mp': 0, 'ms': 0, 'mr': 0, 'vt': 0, 'my': 0, 'mx': 0, 'mz': 0, 'vv': 0, 'vw': 0, 'vx': 0, 'vz': 0, 'fp': 0, 'fq': 0, 'fr': 0, 'fs': 0, 'ft': 0, 'fu': 0, 'fv': 0, 'fw': 0, 'fx': 0, 'fy': 0, 'fz': 0, 'fa': 0, 'fb': 0, 'fc': 0, 'fd': 0, 'fe': 0, 'ff': 0, 'fg': 0, 'fh': 0, 'fi': 12, 'fj': 0, 'fk': 0, 'fl': 1, 'fm': 0, 'fn': 0, 'fo': 0, 'sz': 0, 'sy': 16, 'sx': 0, 'ss': 0, 'sr': 0, 'sq': 0, 'sp': 22, 'sw': 0, 'sv': 0, 'su': 3, 'st': 1, 'sk': 0, 'sj': 0, 'si': 15, 'sh': 5, 'so': 1, 'sn': 0, 'sm': 2, 'sl': 5, 'sc': 0, 'sb': 0, 'sa': 4, 'sg': 0, 'sf': 0, 'se': 9, 'sd': 0, 'lf': 0, 'lg': 1, 'ld': 12, 'le': 20, 'lb': 0, 'lc': 0, 'la': 11, 'ln': 0, 'lo': 1, 'll': 0, 'lm': 0, 'lj': 0, 'lk': 0, 'lh': 0, 'li': 4, 'lv': 9, 'lw': 0, 'lt': 1, 'lu': 3, 'lr': 0, 'ls': 1, 'lp': 3, 'lq': 0, 'lz': 0, 'lx': 0, 'ly': 7, 'wq': 0, 'yh': 0, 'yk': 0, 'yj': 0, 'ym': 0, 'yl': 3, 'yo': 0, 'yn': 0, 'ya': 0, 'yc': 2, 'yb': 1, 'ye': 0, 'yd': 0, 'yg': 1, 'yf': 0, 'yy': 0, 'yx': 0, 'yz': 0, 'yq': 0, 'yp': 2, 'ys': 10, 'yr': 1, 'yu': 0, 'yt': 0, 'yw': 0, 'yv': 0, 'em': 6, 'el': 21, 'eo': 11, 'en': 16, 'ei': 60, 'eh': 0, 'ek': 0, 'ej': 0, 'ee': 0, 'ed': 5, 'eg': 0, 'ef': 0, 'ea': 1, 'ec': 4, 'eb': 0, 'ey': 2, 'ex': 0, 'ez': 0, 'eu': 85, 'et': 0, 'ew': 0, 'ev': 0, 'eq': 0, 'ep': 2, 'es': 5, 'er': 29, 'rt': 2, 'ru': 10, 'rv': 0, 'rw': 0, 'rp': 1, 'rq': 0, 'rr': 0, 'rs': 0, 'rx': 0, 'ry': 2, 'rz': 0, 'rd': 0, 're': 24, 'rf': 0, 'rg': 3, 'ra': 12, 'rb': 0, 'rc': 0, 'rl': 2, 'rm': 0, 'rn': 7, 'ro': 30, 'rh': 0, 'ri': 14, 'rj': 0, 'rk': 2, 'xj': 0, 'xk': 0, 'xh': 0, 'xi': 0, 'xn': 0, 'xo': 0, 'xl': 0, 'xm': 0, 'xb': 0, 'xc': 0, 'xa': 0, 'xf': 0, 'xg': 0, 'xd': 0, 'xe': 0, 'xz': 0, 'xx': 0, 'xy': 0, 'xr': 0, 'xs': 0, 'xp': 1, 'xq': 0, 'xv': 0, 'xw': 0, 'xt': 0, 'xu': 0, 'wy': 8, 'wx': 0, 'kc': 0, 'kb': 0, 'ka': 0, 'kg': 0, 'kf': 0, 'ke': 2, 'kd': 0, 'kk': 0, 'kj': 0, 'ki': 0, 'kh': 0, 'ko': 0, 'kn': 0, 'km': 0, 'kl': 0, 'ks': 0, 'kr': 0, 'kq': 0, 'kp': 0, 'kw': 0, 'kv': 0, 'ku': 0, 'kt': 0, 'kz': 0, 'ky': 0, 'kx': 0, 'dn': 0, 'do': 0, 'dl': 0, 'dm': 0, 'dj': 0, 'dk': 0, 'dh': 0, 'di': 7, 'df': 0, 'dg': 0, 'dd': 0, 'de': 0, 'db': 0, 'dc': 0, 'da': 0, 'dz': 0, 'dx': 0, 'dy': 0, 'dv': 0, 'dw': 0, 'dt': 0, 'du': 2, 'dr': 1, 'ds': 0, 'dp': 0, 'dq': 0, 'qq': 0, 'qp': 0, 'qs': 0, 'qr': 0, 'qu': 0, 'qt': 0, 'qw': 0, 'qv': 0, 'qy': 0, 'qx': 0, 'qz': 0, 'qa': 0, 'qc': 0, 'qb': 0, 'qe': 0, 'qd': 0, 'qg': 0, 'qf': 0, 'qi': 0, 'qh': 0, 'qk': 0, 'qj': 0, 'qm': 0, 'ql': 0, 'qo': 0, 'qn': 0, 'wc': 0, 'wb': 0, 'wa': 0, 'wo': 1, 'wn': 1, 'wm': 1, 'wg': 0, 'wf': 0, 'we': 0, 'wd': 0, 'jx': 0, 'jy': 0, 'jz': 0, 'jt': 0, 'ju': 0, 'jv': 0, 'jw': 0, 'jp': 0, 'jq': 0, 'jr': 0, 'js': 0, 'jl': 0, 'jm': 0, 'jn': 0, 'jo': 0, 'jh': 0, 'ji': 0, 'jj': 0, 'jk': 0, 'jd': 0, 'je': 0, 'jf': 0, 'jg': 0, 'ja': 0, 'jb': 0, 'jc': 0, 'ww': 0, 'wv': 0, 'wu': 0, 'wt': 0, 'ws': 0, 'wr': 0, 'ck': 0, 'cj': 0, 'ci': 85, 'ch': 1, 'co': 13, 'cn': 0, 'cm': 0, 'cl': 15, 'cc': 0, 'cb': 0, 'ca': 0, 'wp': 0, 'cg': 0, 'cf': 0, 'ce': 1, 'cd': 0, 'cz': 0, 'cy': 0, 'cx': 0, 'cs': 3, 'cr': 0, 'cq': 0, 'cp': 0, 'cw': 0, 'cv': 0, 'cu': 7, 'ct': 0, 'pr': 5, 'ps': 3, 'pp': 0, 'pq': 0, 'pv': 0, 'pw': 0, 'pt': 6, 'pu': 0, 'pz': 0, 'px': 0, 'py': 0, 'wz': 0, 'pb': 0, 'pc': 0, 'pa': 17, 'pf': 0, 'pg': 0, 'pd': 0, 'pe': 4, 'pj': 0, 'pk': 0, 'ph': 1, 'pi': 0, 'pn': 0, 'po': 1, 'pl': 0, 'pm': 0, 'iy': 0, 'ix': 0, 'vb': 0, 'iz': 3, 'vd': 0, 've': 1, 'vf': 0, 'vg': 0, 'iq': 1, 'ip': 0, 'is': 42, 'ir': 13, 'iu': 0, 'it': 35, 'iw': 0, 'iv': 6, 'ii': 0, 'ih': 0, 'ik': 0, 'ij': 0, 'im': 0, 'il': 9, 'io': 11, 'in': 5, 'ia': 15, 'vy': 0, 'ic': 31, 'ib': 8, 'ie': 66, 'id': 3, 'ig': 3, 'if': 1, 'wh': 4, 'yi': 0, 'vr': 0, 'vs': 0, 'bd': 0, 'be': 2, 'bf': 0, 'bg': 0, 'ba': 0, 'bb': 0, 'bc': 0, 'bl': 1, 'bm': 1, 'bn': 0, 'bo': 2, 'bh': 0, 'bi': 0, 'bj': 0, 'bk': 0, 'bt': 0, 'bu': 0, 'bv': 0, 'bw': 0, 'bp': 0, 'bq': 0, 'br': 0, 'bs': 2, 'bx': 0, 'by': 0, 'bz': 0, 'oo': 0, 'on': 5, 'om': 0, 'ol': 1, 'ok': 0, 'oj': 0, 'oi': 5, 'oh': 0, 'og': 0, 'of': 0, 'oe': 4, 'od': 0, 'oc': 2, 'ob': 0, 'oa': 5, 'oz': 0, 'oy': 0, 'ox': 1, 'ow': 7, 'ov': 0, 'ou': 0, 'ot': 1, 'os': 1, 'or': 11, 'oq': 0, 'op': 1, 'hz': 0, 'hx': 0, 'hy': 0, 'hr': 0, 'hs': 0, 'hp': 0, 'hq': 0, 'hv': 0, 'hw': 0, 'ht': 10, 'hu': 0, 'hj': 0, 'hk': 0, 'hh': 0, 'hi': 0, 'hn': 0, 'ho': 0, 'hl': 0, 'hm': 0, 'hb': 0, 'hc': 0, 'ha': 12, 'hf': 0, 'hg': 0, 'hd': 0, 'he': 15, 'uy': 0, 'ux': 0, 'uz': 0, 'uu': 0, 'ut': 2, 'uw': 0, 'uv': 0, 'uq': 0, 'up': 2, 'us': 11, 'ur': 11, 'um': 1, 'ul': 2, 'uo': 20, 'un': 0, 'ui': 2, 'uh': 0, 'uk': 0, 'uj': 0, 'ue': 1, 'ud': 1, 'ug': 2, 'uf': 0, 'ua': 22, 'uc': 5, 'ub': 0, 'aa': 0, 'ac': 2, 'ab': 0, 'ae': 1, 'ad': 1, 'ag': 0, 'af': 0, 'ai': 19, 'ah': 0, 'ak': 1, 'aj': 0, 'am': 4, 'al': 14, 'ao': 10, 'an': 25, 'aq': 0, 'ap': 3, 'as': 3, 'ar': 27, 'au': 31, 'at': 5, 'aw': 0, 'av': 0, 'ay': 0, 'ax': 0, 'az': 0, 'nh': 0, 'ni': 1, 'nj': 0, 'nk': 0, 'nl': 0, 'nm': 3, 'nn': 0, 'no': 0, 'na': 15, 'nb': 0, 'nc': 6, 'nd': 2, 'ne': 12, 'nf': 0, 'ng': 8, 'nx': 0, 'ny': 0, 'nz': 0, 'np': 0, 'nq': 0, 'nr': 0, 'ns': 6, 'nt': 4, 'nu': 0, 'nv': 0, 'nw': 0, 'vp': 0, 'vq': 0}
# add_data = {'gw': 1, 'gv': 0, 'gu': 3, 'gt': 2, 'gs': 69, 'gr': 5, 'gq': 1, 'gp': 0, '#y': 1, 'gz': 0, 'gy': 0, 'gx': 0, 'gg': 5, 'gf': 1, 'ge': 5, 'gd': 0, 'gc': 0, 'gb': 0, 'ga': 8, 'go': 1, 'gn': 1, 'gm': 0, 'gl': 2, 'gk': 0, 'gj': 0, 'gi': 8, 'gh': 12, 'tz': 0, 'tx': 0, 'ty': 6, 'tv': 0, 'tw': 5, 'tt': 183, 'tu': 11, 'tr': 54, 'ts': 264, 'tp': 1, 'tq': 0, 'tn': 1, 'to': 23, 'tl': 6, 'tm': 3, 'tj': 1, 'tk': 0, 'th': 24, 'ti': 59, 'tf': 1, 'tg': 10, 'td': 3, 'te': 65, 'tb': 0, 'tc': 0, 'ta': 39, 'vu': 1, 'zl': 0, 'zm': 0, 'zn': 0, 'zo': 0, 'zh': 0, 'zi': 6, 'zj': 0, 'zk': 0, 'zd': 0, 'ze': 5, 'zf': 1, 'zg': 0, 'za': 2, 'zb': 0, 'zc': 0, 'zx': 0, 'zy': 2, 'zz': 4, 'zt': 0, 'zu': 0, 'zv': 0, 'zw': 0, 'zp': 0, 'zq': 0, 'zr': 0, 'zs': 0, '#s': 23, 'wl': 1, '#q': 0, 'va': 0, '#c': 9, 'vc': 0, 'wk': 1, '#p': 10, 'vh': 0, 'wj': 0, 'vi': 10, 'vj': 0, 'vk': 0, 'vl': 1, 'vm': 0, 'wi': 1, '#v': 1, 'vn': 1, 'vo': 0, 'me': 17, 'md': 0, 'mg': 0, 'mf': 0, 'ma': 11, 'mc': 1, 'mb': 1, 'mm': 102, 'ml': 0, 'mo': 7, 'mn': 44, 'mi': 6, 'mh': 1, 'mk': 1, 'mj': 0, 'mu': 2, 'mt': 1, 'mw': 1, 'mv': 0, 'mq': 0, 'mp': 2, 'ms': 47, 'mr': 0, 'vt': 0, 'my': 0, 'mx': 0, 'mz': 0, 'vv': 5, 'vw': 1, '#t': 2, 'vx': 0, 'vz': 0, '#b': 8, 'fp': 0, 'fq': 0, 'fr': 5, 'fs': 23, 'ft': 0, 'fu': 1, 'fv': 0, 'fw': 0, 'fx': 0, 'fy': 1, 'fz': 0, 'fa': 1, 'fb': 0, 'fc': 0, 'fd': 0, 'fe': 2, 'ff': 27, 'fg': 1, 'fh': 0, 'fi': 12, 'fj': 0, 'fk': 0, 'fl': 10, 'fm': 0, 'fn': 0, 'fo': 0, 'sz': 0, 'sy': 7, 'sx': 0, 'ss': 205, 'sr': 1, 'sq': 0, 'sp': 1, 'sw': 1, 'sv': 0, 'su': 7, 'st': 49, 'sk': 2, 'sj': 0, 'si': 101, 'sh': 50, 'so': 3, 'sn': 7, 'sm': 10, 'sl': 2, 'sc': 7, 'sb': 1, 'sa': 13, 'sg': 1, 'sf': 0, 'se': 41, 'sd': 20, 'lf': 0, 'lg': 0, 'ld': 1, 'le': 38, 'lb': 1, 'lc': 0, 'la': 3, 'ln': 0, 'lo': 7, 'll': 128, 'lm': 1, 'lj': 0, 'lk': 2, 'lh': 0, 'li': 79, 'lv': 1, 'lw': 0, 'lt': 7, 'lu': 3, 'lr': 0, 'ls': 97, 'lp': 0, 'lq': 0, 'lz': 0, 'lx': 0, 'ly': 2, 'wq': 1, 'yh': 0, 'yk': 0, 'yj': 0, 'ym': 1, 'yl': 1, 'yo': 0, 'yn': 6, 'ya': 5, 'yc': 2, 'yb': 1, 'ye': 3, 'yd': 0, 'yg': 0, 'yf': 0, 'yy': 2, 'yx': 0, 'yz': 0, 'yq': 0, 'yp': 0, 'ys': 33, 'yr': 1, 'yu': 13, 'yt': 1, 'yw': 1, 'yv': 0, 'em': 6, 'el': 4, 'eo': 5, 'en': 27, 'ei': 4, 'eh': 1, 'ek': 3, 'ej': 0, 'ee': 147, 'ed': 76, 'eg': 0, 'ef': 2, 'ea': 39, 'ec': 8, 'eb': 2, 'ey': 8, 'ex': 2, '#g': 14, 'ez': 0, 'eu': 4, 'et': 6, 'ew': 10, 'ev': 1, 'eq': 0, 'ep': 1, 'es': 417, 'er': 83, 'rt': 29, 'ru': 7, 'rv': 0, 'rw': 1, 'rp': 0, 'rq': 0, 'rr': 132, 'rs': 273, 'rx': 0, 'ry': 10, 'rz': 0, 'rd': 0, 're': 89, 'rf': 1, 'rg': 1, 'ra': 15, 'rb': 2, 'rc': 1, 'rl': 5, 'rm': 9, 'rn': 7, 'ro': 10, 'rh': 2, 'ri': 64, 'rj': 0, 'rk': 0, 'xj': 0, 'xk': 0, 'xh': 6, 'xi': 1, 'xn': 0, 'xo': 3, 'xl': 0, 'xm': 1, 'xb': 0, 'xc': 18, 'xa': 0, 'xf': 0, 'xg': 0, 'xd': 0, 'xe': 1, 'xz': 0, 'xx': 1, 'xy': 0, 'xr': 0, 'xs': 2, 'xp': 0, 'xq': 0, 'xv': 0, 'xw': 0, 'xt': 0, 'xu': 0, 'wy': 0, 'wx': 0, '#d': 8, 'kc': 0, 'kb': 4, 'ka': 2, 'kg': 0, 'kf': 0, 'ke': 9, 'kd': 1, 'kk': 1, 'kj': 0, 'ki': 1, 'kh': 1, 'ko': 2, 'kn': 0, 'km': 0, 'kl': 1, 'ks': 95, 'kr': 0, 'kq': 0, 'kp': 1, 'kw': 0, 'kv': 0, 'ku': 1, 'kt': 0, 'kz': 0, 'ky': 4, 'kx': 0, 'dn': 9, 'do': 13, 'dl': 6, 'dm': 1, 'dj': 0, 'dk': 0, 'dh': 0, 'di': 9, 'df': 2, 'dg': 0, 'dd': 17, 'de': 14, 'db': 0, 'dc': 3, 'da': 18, 'dz': 0, 'dx': 0, 'dy': 5, 'dv': 0, 'dw': 0, 'dt': 0, 'du': 0, 'dr': 6, 'ds': 119, 'dp': 0, 'dq': 0, 'qq': 0, 'qp': 0, 'qs': 0, 'qr': 0, 'qu': 1, 'qt': 0, 'qw': 0, 'qv': 0, 'qy': 0, 'qx': 0, 'qz': 0, 'qa': 0, 'qc': 0, 'qb': 0, 'qe': 0, 'qd': 0, 'qg': 0, 'qf': 0, 'qi': 1, 'qh': 0, 'qk': 0, 'qj': 0, 'qm': 0, 'ql': 0, 'qo': 0, 'qn': 0, '#k': 17, '#j': 1, '#e': 26, '#i': 5, '#h': 3, 'wc': 0, 'wb': 0, 'wa': 0, 'wo': 0, 'wn': 2, 'wm': 0, 'wg': 0, 'wf': 0, 'we': 10, 'wd': 1, 'jx': 0, 'jy': 0, 'jz': 0, '#l': 5, 'jt': 0, 'ju': 1, 'jv': 0, 'jw': 0, 'jp': 0, 'jq': 0, 'jr': 0, 'js': 0, 'jl': 0, 'jm': 0, 'jn': 0, 'jo': 0, 'jh': 0, 'ji': 0, 'jj': 0, 'jk': 0, 'jd': 0, 'je': 0, 'jf': 0, 'jg': 0, '#w': 2, 'ja': 0, 'jb': 0, 'jc': 0, 'ww': 4, 'wv': 0, 'wu': 2, 'wt': 0, 'ws': 8, 'wr': 1, 'ck': 3, 'cj': 0, 'ci': 50, 'ch': 18, 'co': 7, 'cn': 1, 'cm': 1, 'cl': 1, 'cc': 54, 'cb': 0, 'ca': 19, 'wp': 0, 'cg': 0, 'cf': 0, 'ce': 13, 'cd': 1, 'cz': 0, 'cy': 0, 'cx': 1, '#r': 6, 'cs': 25, 'cr': 7, 'cq': 0, 'cp': 1, 'cw': 0, 'cv': 4, 'cu': 8, 'ct': 7, 'pr': 29, 'ps': 52, 'pp': 70, 'pq': 0, 'pv': 1, 'pw': 1, 'pt': 9, 'pu': 1, 'pz': 0, 'px': 0, 'py': 0, 'wz': 0, 'pb': 0, 'pc': 1, 'pa': 23, 'pf': 0, 'pg': 0, 'pd': 1, 'pe': 10, 'pj': 0, 'pk': 0, 'ph': 20, 'pi': 3, 'pn': 0, 'po': 26, 'pl': 2, 'pm': 0, 'iy': 0, 'ix': 1, 'vb': 2, 'iz': 1, 'vd': 0, 've': 36, 'vf': 0, 'vg': 0, 'iq': 0, 'ip': 1, 'is': 30, 'ir': 9, 'iu': 11, 'it': 29, 'iw': 0, 'iv': 0, 'ii': 69, 'ih': 1, 'ik': 1, 'ij': 2, 'im': 11, 'il': 17, 'io': 27, 'in': 33, 'ia': 10, 'vy': 0, 'ic': 13, 'ib': 3, 'ie': 25, 'id': 13, 'ig': 1, 'if': 0, '#x': 1, 'wh': 1, 'yi': 2, '#u': 11, 'vr': 0, '#f': 11, '#o': 2, '#n': 2, '#m': 6, 'vs': 0, 'bd': 0, 'be': 7, 'bf': 0, 'bg': 1, 'ba': 3, 'bb': 11, 'bc': 0, 'bl': 15, 'bm': 0, 'bn': 1, 'bo': 1, 'bh': 0, 'bi': 50, 'bj': 0, 'bk': 0, 'bt': 0, 'bu': 0, 'bv': 3, 'bw': 0, 'bp': 0, 'bq': 0, 'br': 5, 'bs': 16, 'bx': 0, 'by': 0, 'bz': 0, 'oo': 64, 'on': 13, 'om': 3, 'ol': 6, 'ok': 0, 'oj': 1, 'oi': 28, 'oh': 0, 'og': 1, 'of': 2, 'oe': 7, 'od': 3, 'oc': 1, 'ob': 1, 'oa': 14, 'oz': 1, 'oy': 1, 'ox': 0, 'ow': 0, 'ov': 1, 'ou': 19, 'ot': 4, 'os': 59, 'or': 16, 'oq': 0, 'op': 30, '#a': 46, 'hz': 0, 'hx': 0, 'hy': 3, 'hr': 16, 'hs': 24, 'hp': 0, 'hq': 0, 'hv': 0, 'hw': 5, 'ht': 22, 'hu': 1, 'hj': 2, 'hk': 0, 'hh': 18, 'hi': 17, 'hn': 1, 'ho': 4, 'hl': 1, 'hm': 0, 'hb': 1, 'hc': 0, 'ha': 4, 'hf': 0, 'hg': 10, 'hd': 1, 'he': 24, 'uy': 3, 'ux': 2, 'uz': 0, 'uu': 26, 'ut': 27, 'uw': 0, 'uv': 0, 'uq': 0, 'up': 3, 'us': 19, 'ur': 49, 'um': 3, 'ul': 3, 'uo': 1, 'un': 9, 'ui': 24, 'uh': 1, 'uk': 1, 'uj': 1, 'ue': 9, 'ud': 0, 'ug': 0, 'uf': 0, 'ua': 15, 'uc': 3, 'ub': 0, 'aa': 15, 'ac': 14, 'ab': 1, 'ae': 10, 'ad': 7, 'ag': 1, 'af': 0, 'ai': 33, 'ah': 1, 'ak': 4, 'aj': 1, 'am': 2, 'al': 31, 'ao': 12, 'an': 39, 'aq': 3, 'ap': 4, 'as': 134, 'ar': 28, 'au': 28, 'at': 7, 'aw': 1, 'av': 0, 'ay': 4, 'ax': 1, 'az': 1, 'nh': 0, 'ni': 34, 'nj': 0, 'nk': 1, 'nl': 1, 'nm': 26, 'nn': 99, 'no': 12, 'na': 15, 'nb': 5, 'nc': 7, 'nd': 13, 'ne': 52, 'nf': 4, 'ng': 17, 'nx': 0, 'ny': 1, 'nz': 0, 'np': 0, 'nq': 0, 'nr': 2, 'ns': 156, 'nt': 53, 'nu': 1, 'nv': 1, 'nw': 0, 'vp': 1, '#z': 2, 'vq': 0}
# del_data = {'tz': 0, 'tx': 0, 'ty': 2, 'tv': 0, 'tw': 4, 'tt': 137, 'tu': 14, 'tr': 203, 'ts': 5, 'tp': 1, 'tq': 0, 'tn': 3, 'to': 11, 'tl': 31, 'tm': 3, 'tj': 0, 'tk': 0, 'th': 49, 'ti': 427, 'tf': 1, 'tg': 7, 'td': 0, 'te': 76, 'tb': 1, 'tc': 2, 'ta': 24, 'me': 33, 'md': 0, 'mg': 0, 'mf': 0, 'ma': 15, 'mc': 0, 'mb': 10, 'mm': 180, 'ml': 0, 'mo': 7, 'mn': 7, 'mi': 42, 'mh': 1, 'mk': 0, 'mj': 0, 'mu': 4, 'mt': 0, 'mw': 0, 'mv': 0, 'mq': 0, 'mp': 31, 'ms': 9, 'mr': 0, 'my': 0, 'mx': 0, 'mz': 0, 'fp': 0, 'fq': 0, 'fr': 11, 'fs': 0, 'ft': 8, 'fu': 1, 'fv': 0, 'fw': 0, 'fx': 0, 'fy': 1, 'fz': 0, 'fa': 4, 'fb': 0, 'fc': 0, 'fd': 0, 'fe': 13, 'ff': 46, 'fg': 0, 'fh': 0, 'fi': 79, 'fj': 0, 'fk': 0, 'fl': 12, 'fm': 0, 'fn': 0, 'fo': 4, 'yi': 1, 'yh': 0, 'yk': 0, 'yj': 0, 'ym': 2, 'yl': 1, 'yo': 1, 'yn': 1, 'ya': 2, 'yc': 34, 'yb': 1, 'ye': 2, 'yd': 0, 'yg': 1, 'yf': 0, 'yy': 0, 'yx': 0, 'yz': 0, 'yq': 0, 'yp': 1, 'ys': 17, 'yr': 0, 'yu': 0, 'yt': 1, 'yw': 1, 'yv': 0, 'rt': 68, 'ru': 0, 'rv': 10, 'rw': 1, 'rp': 2, 'rq': 0, 'rr': 277, 'rs': 103, 'rx': 0, 'ry': 27, 'rz': 0, 'rd': 19, 're': 188, 'rf': 0, 'rg': 11, 'ra': 63, 'rb': 4, 'rc': 12, 'rl': 33, 'rm': 7, 'rn': 157, 'ro': 21, 'rh': 5, 'ri': 132, 'rj': 0, 'rk': 3, 'kc': 0, 'kb': 0, 'ka': 4, 'kg': 8, 'kf': 1, 'ke': 15, 'kd': 1, 'kk': 1, 'kj': 0, 'ki': 5, 'kh': 1, 'ko': 0, 'kn': 17, 'km': 0, 'kl': 3, 'ks': 5, 'kr': 1, 'kq': 0, 'kp': 0, 'kw': 1, 'kv': 0, 'ku': 0, 'kt': 0, 'kz': 0, 'ky': 0, 'kx': 0, 'dn': 3, 'do': 3, 'dl': 8, 'dm': 4, 'dj': 1, 'dk': 1, 'dh': 0, 'di': 62, 'df': 0, 'dg': 10, 'dd': 25, 'de': 45, 'db': 0, 'dc': 7, 'da': 12, 'dz': 0, 'dx': 0, 'dy': 6, 'dv': 2, 'dw': 0, 'dt': 0, 'du': 3, 'dr': 11, 'ds': 1, 'dp': 0, 'dq': 0, 'wg': 0, 'wf': 1, 'we': 11, 'wd': 1, 'wc': 0, 'wb': 0, 'wa': 40, 'wo': 2, 'wn': 2, 'wm': 0, 'wl': 1, 'wk': 0, 'wj': 0, 'wi': 15, 'wh': 11, 'ww': 0, 'wv': 0, 'wu': 0, 'wt': 0, 'ws': 24, 'wr': 2, 'wq': 0, 'wp': 0, 'wz': 0, 'wy': 0, 'wx': 0, 'pr': 58, 'ps': 1, 'pp': 93, 'pq': 0, 'pv': 0, 'pw': 0, 'pt': 18, 'pu': 2, 'pz': 0, 'px': 0, 'py': 0, 'pb': 0, 'pc': 0, 'pa': 25, 'pf': 0, 'pg': 0, 'pd': 0, 'pe': 22, 'pj': 0, 'pk': 0, 'ph': 12, 'pi': 15, 'pn': 0, 'po': 30, 'pl': 28, 'pm': 1, 'iy': 1, 'ix': 0, 'iz': 7, 'iq': 0, 'ip': 7, 'is': 71, 'ir': 16, 'iu': 1, 'it': 64, 'iw': 0, 'iv': 1, 'ii': 1, 'ih': 0, 'ik': 0, 'ij': 0, 'im': 14, 'il': 38, 'io': 41, 'in': 82, 'ia': 26, 'ic': 60, 'ib': 1, 'ie': 23, 'id': 26, 'ig': 9, 'if': 1, 'bd': 0, 'be': 22, 'bf': 0, 'bg': 0, 'ba': 2, 'bb': 2, 'bc': 1, 'bl': 26, 'bm': 0, 'bn': 0, 'bo': 2, 'bh': 0, 'bi': 183, 'bj': 0, 'bk': 0, 'bt': 0, 'bu': 6, 'bv': 1, 'bw': 0, 'bp': 0, 'bq': 0, 'br': 6, 'bs': 17, 'bx': 0, 'by': 0, 'bz': 0, '#g': 7, 'uy': 1, 'ux': 0, 'uz': 0, 'uu': 0, 'ut': 66, 'uw': 0, 'uv': 0, 'uq': 0, 'up': 0, 'us': 31, 'ur': 129, 'um': 2, 'ul': 39, 'uo': 1, 'un': 111, 'ui': 28, 'uh': 0, 'uk': 0, 'uj': 0, 'ue': 15, 'ud': 10, 'ug': 1, 'uf': 0, 'ua': 26, 'uc': 9, 'ub': 6, '#e': 20, 'nh': 0, 'ni': 191, 'nj': 0, 'nk': 0, 'nl': 0, 'nm': 17, 'nn': 144, 'no': 21, 'na': 21, 'nb': 0, 'nc': 42, 'nd': 71, 'ne': 68, 'nf': 1, 'ng': 160, 'nx': 0, 'ny': 2, 'nz': 0, 'np': 0, 'nq': 0, 'nr': 0, 'ns': 127, 'nt': 87, 'nu': 43, 'nv': 1, 'nw': 1, 'gw': 0, 'gv': 0, 'gu': 22, 'gt': 1, 'gs': 7, 'gr': 52, 'gq': 0, 'gp': 0, 'gz': 0, 'gy': 1, 'gx': 0, 'gg': 37, 'gf': 1, 'ge': 83, 'gd': 2, 'gc': 0, 'gb': 0, 'ga': 25, 'go': 4, 'gn': 29, 'gm': 0, 'gl': 3, 'gk': 0, 'gj': 0, 'gi': 39, 'gh': 25, 'zl': 0, 'zm': 0, 'zn': 0, 'zo': 0, 'zh': 0, 'zi': 0, 'zj': 0, 'zk': 0, 'zd': 0, 'ze': 2, 'zf': 0, 'zg': 0, 'za': 1, 'zb': 0, 'zc': 0, 'zx': 0, 'zy': 0, 'zz': 2, 'zt': 0, 'zu': 0, 'zv': 0, 'zw': 0, 'zp': 0, 'zq': 0, 'zr': 0, 'zs': 0, '#s': 26, '#q': 0, 'sz': 0, 'sy': 1, 'sx': 0, 'ss': 265, 'sr': 4, 'sq': 0, 'sp': 30, 'sw': 0, 'sv': 0, 'su': 21, 'st': 124, 'sk': 0, 'sj': 0, 'si': 231, 'sh': 18, 'so': 30, 'sn': 0, 'sm': 1, 'sl': 2, 'sc': 27, 'sb': 0, 'sa': 16, 'sg': 0, 'sf': 1, 'se': 74, 'sd': 0, 'lf': 0, 'lg': 0, 'ld': 6, 'le': 48, 'lb': 0, 'lc': 1, 'la': 24, 'ln': 0, 'lo': 29, 'll': 211, 'lm': 2, 'lj': 0, 'lk': 0, 'lh': 0, 'li': 217, 'lv': 2, 'lw': 0, 'lt': 7, 'lu': 3, 'lr': 2, 'ls': 12, 'lp': 0, 'lq': 0, 'lz': 0, 'lx': 0, 'ly': 11, 'em': 9, 'el': 32, 'eo': 19, 'en': 76, 'ei': 6, 'eh': 1, 'ek': 0, 'ej': 0, 'ee': 89, 'ed': 74, 'eg': 1, 'ef': 3, 'ea': 80, 'ec': 50, 'eb': 1, 'ey': 1, 'ex': 7, 'ez': 0, 'eu': 8, 'et': 34, 'ew': 1, 'ev': 2, 'eq': 1, 'ep': 9, 'es': 223, 'er': 237, 'xj': 0, 'xk': 0, 'xh': 1, 'xi': 0, 'xn': 0, 'xo': 0, 'xl': 0, 'xm': 0, 'xb': 0, 'xc': 17, 'xa': 1, 'xf': 0, 'xg': 0, 'xd': 0, 'xe': 3, 'xz': 0, 'xx': 0, 'xy': 1, 'xr': 0, 'xs': 0, 'xp': 6, 'xq': 0, 'xv': 0, 'xw': 0, 'xt': 5, 'xu': 0, '#f': 20, 'qq': 0, 'qp': 0, 'qs': 0, 'qr': 0, 'qu': 18, 'qt': 0, 'qw': 0, 'qv': 0, 'qy': 0, 'qx': 0, 'qz': 0, 'qa': 0, 'qc': 0, 'qb': 0, 'qe': 0, 'qd': 0, 'qg': 0, 'qf': 0, 'qi': 0, 'qh': 0, 'qk': 0, 'qj': 0, 'qm': 0, 'ql': 0, 'qo': 0, 'qn': 0, '#o': 5, '#n': 5, '#m': 16, '#c': 41, '#b': 14, '#a': 20, '#k': 6, '#j': 3, '#i': 20, '#h': 6, 'jx': 0, 'jy': 0, 'jz': 0, '#l': 22, 'jt': 0, 'ju': 1, 'jv': 0, 'jw': 0, 'jp': 0, 'jq': 0, 'jr': 0, 'js': 0, 'jl': 0, 'jm': 0, 'jn': 1, 'jo': 1, 'jh': 0, 'ji': 0, 'jj': 0, 'jk': 0, 'jd': 0, 'je': 1, 'jf': 0, 'jg': 0, '#w': 24, 'ja': 0, 'jb': 0, 'jc': 0, '#z': 2, '#y': 0, '#x': 0, 'ck': 9, 'cj': 0, 'ci': 320, 'ch': 24, 'co': 33, 'cn': 0, 'cm': 0, 'cl': 17, 'cc': 70, 'cb': 0, 'ca': 37, 'cg': 0, 'cf': 0, 'ce': 63, 'cd': 0, 'cz': 0, 'cy': 1, 'cx': 0, '#r': 28, 'cs': 6, 'cr': 46, 'cq': 0, 'cp': 0, 'cw': 0, 'cv': 0, 'cu': 17, 'ct': 54, '#d': 31, '#p': 17, '#v': 1, '#u': 2, '#t': 6, 'va': 9, 'vb': 0, 'vc': 0, 'vd': 0, 've': 58, 'vf': 0, 'vg': 0, 'vh': 0, 'vi': 31, 'vj': 0, 'vk': 0, 'vl': 0, 'vm': 0, 'vn': 0, 'vo': 2, 'vp': 0, 'vq': 0, 'vr': 1, 'vs': 0, 'vt': 0, 'vu': 0, 'vv': 0, 'vw': 0, 'vx': 0, 'vy': 1, 'vz': 0, 'oo': 26, 'on': 70, 'om': 9, 'ol': 13, 'ok': 0, 'oj': 1, 'oi': 4, 'oh': 0, 'og': 5, 'of': 0, 'oe': 8, 'od': 6, 'oc': 3, 'ob': 4, 'oa': 11, 'oz': 0, 'oy': 1, 'ox': 0, 'ow': 5, 'ov': 2, 'ou': 47, 'ot': 13, 'os': 20, 'or': 98, 'oq': 0, 'op': 20, 'hz': 0, 'hx': 0, 'hy': 1, 'hr': 15, 'hs': 1, 'hp': 0, 'hq': 0, 'hv': 0, 'hw': 1, 'ht': 26, 'hu': 0, 'hj': 0, 'hk': 0, 'hh': 25, 'hi': 24, 'hn': 9, 'ho': 22, 'hl': 7, 'hm': 1, 'hb': 12, 'hc': 1, 'ha': 15, 'hf': 0, 'hg': 0, 'hd': 3, 'he': 20, 'aa': 0, 'ac': 58, 'ab': 7, 'ae': 3, 'ad': 21, 'ag': 18, 'af': 5, 'ai': 61, 'ah': 8, 'ak': 4, 'aj': 0, 'am': 5, 'al': 43, 'ao': 0, 'an': 53, 'aq': 0, 'ap': 9, 'as': 28, 'ar': 98, 'au': 62, 'at': 53, 'aw': 0, 'av': 1, 'ay': 2, 'ax': 0, 'az': 0}

def check_word_exist(word):
    for item in wn.synsets(word):
        if str(item)[8:].split('.')[0] == word:
            return 1
    return 0

def get_definition(word):
    data = wn.synsets(word)
    for item in wn.synsets(word):
        if str(item)[8:].split('.')[0] == word:
            print item.definition()

def check_spell(word):
    alternatives = []
    for item in all_words:
        if damerau_levenshtein_distance(word,item) == 1:
            alternatives.append(item)
    return alternatives

def checking(word):
    if check_word_exist(word) == 0:
        print "WORD NOT FOUND, ALTERNATIVES\n"
        alternatives = check_spell(word)
        print alternatives
    else:
        print "WORD FOUND, DEFINITIONS\n"
        get_definition(word)

def get_diff(word,item):
    for index,char in enumerate(word):
        if index == len(word)-1 or char != item[index]:
            return index

def get_score(word,alternatives,sub_data,int_data,add_data,del_data):
    result = []
    for item in alternatives:
        val = get_diff(word,item)
        if len(word) - len(item) == 0:
            if editdistance.eval(word,item) == 1:
                value = sub_data[word[val]+item[val]]
                print "substitute",word[val],item[val],value
            else:
                value = int_data[word[val]+item[val]]
                print "interchange",word[val],item[val],value
        if len(word) - len(item) == -1:
            if val == 0:
                value = add_data['#'+item[val]]
            else:
                value = add_data[item[val-1]+item[val]]
            print "added",item[val],value
        if len(word) - len(item) == 1:
            if val == 0:
                value = del_data['#'+word[val]]
            else:
                value = del_data[word[val-1]+word[val]]  
            print "deleted",word[val],value
        result.append(value)
    return result

# all_words = wn.all_lemma_names()
# word = 'acress'
# alternatives = check_spell(word)
# print alternatives

# result = get_score(word,alternatives,sub_data,int_data,add_data,del_data)
# print alternatives[result.index(max(result))]

# file = '1_1_all_fullalpha.txt'
# f = open(file)
# count,val = 0,[]
# for item in f.readlines():
#     count += 1
#     array = item.strip().split('\t')
#     val.append(int(array[3]))

def pandas_helper():
    s = pd.Series([1,3,5,np.nan,6,8])
    dates = pd.date_range('20130101', periods=6)
    df = pd.DataFrame(np.random.randn(6,4), index=dates, columns=list('ABCD'))
    df2 = pd.DataFrame({'A':1.,'B':pd.Timestamp('20130102'),'C':pd.Series(1,index=list(range(4)),dtype='float32'),'D':np.array([3]*4,dtype='int32'),'E':pd.Categorical(["test","train","test","train"]),'F':'foo'})
    print df2.dtypes
    print df.head()
    print df.tail(3)
    print df.index
    print df.columns
    print df.values
    print df.describe()
    print df.T
    print df.sort_index(axis=1, ascending=False)
    print df.sort_values(by='B')
    print df['A']
    print df[0:3]
    print df['20130102':'20130104']
    print df.loc[dates[0]]
    print df.loc[:,['A','B']]
    print df.loc['20130102':'20130104',['A','B']]
    print df.loc['20130102',['A','B']]
    print df.loc[dates[0],'A']
    print df.at[dates[0],'A']
    print df.iloc[3]
    print df.iloc[3:5,0:2]
    print df.iloc[[1,2,4],[0,2]]
    print df.iloc[1:3,:]
    print df.iloc[:,1:3]
    print df.iloc[1,1]
    print df.iat[1,1]
    print df[df.A > 0]
    print df[df > 0]
    df2 = df.copy()
    df2['E'] = ['one', 'one','two','three','four','three']
    print df2
    print df2[df2['E'].isin(['two','four'])]
    s1 = pd.Series([1,2,3,4,5,6], index=pd.date_range('20130102', periods=6))
    print s1
    df['F'] = s1
    df.at[dates[0],'A'] = 0
    df.iat[0,1] = 0
    df.loc[:,'D'] = np.array([5] * len(df))
    print df
    df2 = df.copy()
    df2[df2 > 0] = -df2
    print df2
    df1 = df.reindex(index=dates[0:4], columns=list(df.columns) + ['E'])
    df1.loc[dates[0]:dates[1],'E'] = 1
    print df1
    print df1.dropna(how='any')
    print df1.fillna(value=5)
    print pd.isnull(df1)
    print df.mean()
    print df.mean(1)
    s = pd.Series([1,3,5,np.nan,6,8], index=dates).shift(2)
    print s
    print df.apply(np.cumsum)
    print df.apply(lambda x: x.max() - x.min())
    s = pd.Series(np.random.randint(0, 7, size=10))
    print s
    print s.value_counts()
    s = pd.Series(['A', 'B', 'C', 'Aaba', 'Baca', np.nan, 'CABA', 'dog', 'cat'])
    print s.str.lower()
    df = pd.DataFrame(np.random.randn(10, 4))
    print df
    pieces = [df[:3], df[3:7], df[7:]]
    print pd.concat(pieces)
    left = pd.DataFrame({'key': ['foo', 'foo'], 'lval': [1, 2]})
    right = pd.DataFrame({'key': ['foo', 'foo'], 'rval': [4, 5]})
    print left
    print right
    pd.merge(left, right, on='key')
    left = pd.DataFrame({'key': ['foo', 'bar'], 'lval': [1, 2]})
    right = pd.DataFrame({'key': ['foo', 'bar'], 'rval': [4, 5]})
    print left
    print right
    pd.merge(left, right, on='key')
    df = pd.DataFrame(np.random.randn(8, 4), columns=['A','B','C','D'])
    print df
    s = df.iloc[3]
    print df.append(s, ignore_index=True)
    df = pd.DataFrame({'A' : ['foo', 'bar', 'foo', 'bar', 'foo', 'bar', 'foo', 'foo'], 'B' : ['one', 'one', 'two', 'three', 'two', 'two', 'one', 'three'], 'C' : np.random.randn(8), 'D' : np.random.randn(8)})
    print df
    print df.groupby('A').sum()
    print df.groupby(['A','B']).sum()
    tuples = list(zip(*[['bar', 'bar', 'baz', 'baz', 'foo', 'foo', 'qux', 'qux'], ['one', 'two', 'one', 'two', 'one', 'two', 'one', 'two']]))
    index = pd.MultiIndex.from_tuples(tuples, names=['first', 'second'])
    df = pd.DataFrame(np.random.randn(8, 2), index=index, columns=['A', 'B'])
    df2 = df[:4]
    print df2
    stacked = df2.stack()
    print stacked
    print stacked.unstack()
    print stacked.unstack(1)
    print stacked.unstack(0)
    df = pd.DataFrame({'A' : ['one', 'one', 'two', 'three'] * 3, 'B' : ['A', 'B', 'C'] * 4, 'C' : ['foo', 'foo', 'foo', 'bar', 'bar', 'bar'] * 2, 'D' : np.random.randn(12), 'E' : np.random.randn(12)})
    print df
    print pd.pivot_table(df, values='D', index=['A', 'B'], columns=['C'])
    rng = pd.date_range('1/1/2012', periods=100, freq='S')
    ts = pd.Series(np.random.randint(0, 500, len(rng)), index=rng)
    print ts.resample('5Min').sum()
    rng = pd.date_range('3/6/2012 00:00', periods=5, freq='D')
    ts = pd.Series(np.random.randn(len(rng)), rng)
    print ts
    ts_utc = ts.tz_localize('UTC')
    print ts_utc
    print ts_utc.tz_convert('US/Eastern')
    rng = pd.date_range('1/1/2012', periods=5, freq='M')
    ts = pd.Series(np.random.randn(len(rng)), index=rng)
    print ts
    ps = ts.to_period()
    print ps
    ps.to_timestamp()
    prng = pd.period_range('1990Q1', '2000Q4', freq='Q-NOV')
    ts = pd.Series(np.random.randn(len(prng)), prng)
    ts.index = (prng.asfreq('M', 'e') + 1).asfreq('H', 's') + 9
    print ts.head()
    df = pd.DataFrame({"id":[1,2,3,4,5,6], "raw_grade":['a', 'b', 'b', 'a', 'a', 'e']})
    df["grade"] = df["raw_grade"].astype("category")
    print df["grade"]
    df["grade"].cat.categories = ["very good", "good", "very bad"]
    df["grade"] = df["grade"].cat.set_categories(["very bad", "bad", "medium", "good", "very good"])
    print df["grade"]
    print df.sort_values(by="grade")
    print df.groupby("grade").size()
    ts = pd.Series(np.random.randn(1000), index=pd.date_range('1/1/2000', periods=1000))
    ts = ts.cumsum()
    ts.plot()
    df = pd.DataFrame(np.random.randn(1000, 4), index=ts.index, columns=['A', 'B', 'C', 'D'])
    df = df.cumsum()
    plt.figure(); df.plot(); plt.legend(loc='best')
    df.to_csv('foo.csv')
    print pd.read_csv('foo.csv')
    # df.to_hdf('foo.h5','df')
    # print pd.read_hdf('foo.h5','df')
    # df.to_excel('foo.xlsx', sheet_name='Sheet1')
    # print pd.read_excel('foo.xlsx', 'Sheet1', index_col=None, na_values=['NA'])

# pandas_helper()

# sudoku = [
#     [4,0,0,0,0,0,8,0,5],
#     [0,3,0,0,0,0,0,0,0],
#     [0,0,0,7,0,0,0,0,0],
#     [0,2,0,0,0,0,0,6,0],
#     [0,0,0,0,8,0,4,0,0],
#     [0,0,0,0,1,0,0,0,0],
#     [0,0,0,6,0,3,0,7,0],
#     [5,0,0,2,0,0,0,0,0],
#     [1,0,4,0,0,0,0,0,0]]
# sudoku = [
#     [7,3,0,9,0,0,8,0,0],
#     [5,1,2,0,0,0,0,0,0],
#     [0,0,0,7,0,0,0,2,3],
#     [3,9,0,1,2,0,0,0,5],
#     [2,0,1,0,0,0,3,0,8],
#     [6,0,0,0,3,8,0,1,2],
#     [4,6,3,2,7,1,5,8,9],
#     [0,0,0,0,0,0,2,3,6],
#     [0,0,5,0,0,3,0,7,0]]
# sudoku = [
#     [2,0,0,8,0,4,0,0,6],
#     [0,0,6,0,0,0,5,0,0],
#     [0,7,4,0,0,0,9,2,0],
#     [3,0,0,0,4,0,0,0,7],
#     [0,0,0,3,0,5,0,0,0],
#     [4,0,0,0,6,0,0,0,9],
#     [0,1,9,0,0,0,7,4,0],
#     [0,0,8,0,0,0,2,0,0],
#     [5,0,0,6,0,8,0,0,1]]
# sudoku = [
#     [0,3,0,9,0,0,8,0,0],
#     [5,1,2,0,0,0,0,0,0],
#     [0,0,0,7,0,0,0,2,3],
#     [3,9,0,1,2,0,0,0,5],
#     [2,0,1,0,0,0,3,0,8],
#     [6,0,0,0,3,8,0,1,2],
#     [4,6,0,0,0,1,0,0,0],
#     [0,0,0,0,0,0,2,3,6],
#     [0,0,5,0,0,3,0,7,0]]
# sudoku = [
#     [0,9,3,0,0,0,5,0,0],
#     [6,0,0,7,1,0,0,0,0],
#     [0,0,0,0,0,4,0,0,3],
#     [0,8,6,0,0,0,1,0,0],
#     [1,0,0,6,4,8,0,0,2],
#     [0,0,2,0,0,0,6,8,0],
#     [3,0,0,1,0,0,0,0,0],
#     [0,0,0,0,5,2,0,0,7],
#     [0,0,4,0,0,0,8,1,0]]
# sudoku = [
#     [0,1,2,0,0,0,0,9,4],
#     [0,0,0,4,0,0,2,7,0],
#     [0,0,0,6,0,0,0,0,0],
#     [0,0,0,0,6,0,7,8,0],
#     [3,0,0,1,0,4,0,0,9],
#     [0,2,5,0,9,0,0,0,0],
#     [0,0,0,0,0,5,0,0,0],
#     [0,3,1,0,0,8,0,0,0],
#     [9,7,0,0,0,0,4,5,0]]

def elements_possible(sudoku,row_num,col_num):
    row_elem = sudoku[row_num]
    col_elem = [sudoku[row][col_num] for row in range(9)]
    sq_elem = []
    sq_row = row_num-row_num%3
    sq_col = col_num-col_num%3
    for irow in range(sq_row,sq_row+3):
        for icol in range(sq_col,sq_col+3):
            sq_elem.append(sudoku[irow][icol])
    x = list(set(row_elem+col_elem+sq_elem))
    p = [0,1,2,3,4,5,6,7,8,9]
    for elem in x: p.remove(elem)
    return p

def check_unfilled_cell(sudoku):
    count = 0
    for irow in range(9):
        for icol in range(9):
            if sudoku[irow][icol] == 0: count += 1
    return count

def check_cell_to_number(sudoku):
    flag = 1
    while flag == 1:
        elem_poss = {}
        flag = 0
        for ii in range(9):
            for jj in range(9):
                if sudoku[ii][jj] == 0:
                    item = elements_possible(sudoku,ii,jj)
                    if len(item) == 1:
                        # print "cell number "+str(item[0])+" row "+str(ii)+" col "+str(jj)
                        sudoku[ii][jj] = item[0]
                        flag = 1
                    else: elem_poss[str(ii)+','+str(jj)] = item    
    return sudoku,elem_poss

def check_number_to_row_col_sq(elem_poss,sudoku):
    for number in range(1,10):
        number_rows = []
        number_cols = []
        number_sqrs = []
        for cell in elem_poss.keys():
            if number in elem_poss[cell]:
                rownum = int(cell.split(',')[0])
                colnum = int(cell.split(',')[1])
                number_rows.append(rownum)
                number_cols.append(colnum)
                number_sqrs.append((rownum/3)*3+colnum/3)
        for item in Counter(number_rows):
            if Counter(number_rows)[item] == 1: 
                for cell in elem_poss.keys():
                    rownum = int(cell.split(',')[0])
                    colnum = int(cell.split(',')[1])
                    if rownum == item and number in elem_poss[cell]:
                        # print "row number "+str(number)+" row "+str(item)+" col "+str(colnum)
                        sudoku[item][colnum] = number

        for item in Counter(number_cols):
            if Counter(number_cols)[item] == 1: 
                for cell in elem_poss.keys():
                    rownum = int(cell.split(',')[0])
                    colnum = int(cell.split(',')[1])
                    if colnum == item and number in elem_poss[cell]: 
                        # print "col number "+str(number)+" row "+str(rownum)+" col "+str(item)
                        sudoku[rownum][item] = number

        for item in Counter(number_sqrs):
            if Counter(number_sqrs)[item] == 1:
                for cell in elem_poss.keys():
                    rownum = int(cell.split(',')[0])
                    colnum = int(cell.split(',')[1])
                    if rownum in range((item/3)*3,(item/3)*3+3) and colnum in range((item%3)*3,(item%3)*3+3) and number in elem_poss[cell]: 
                        # print "sq number "+str(number)+" row "+str(rownum)+" col "+str(colnum)
                        sudoku[rownum][colnum] = number

    return sudoku

def check_same(feature,elem_poss):
    for index in range(9):
        temp = []
        row,col,sqr = [],[],[]
        for item in elem_poss.keys():
            if feature == 'row': val = int(item.split(',')[0])
            if feature == 'col': val = int(item.split(',')[1])
            if feature == 'sqr': val = (int(item.split(',')[0])/3)*3+(int(item.split(',')[1])/3)
            if val == index:
                temp.append(elem_poss[item])
                row.append(int(item.split(',')[0]))
                col.append(int(item.split(',')[1]))
                sqr.append((int(item.split(',')[0])/3)*3+(int(item.split(',')[1])/3))
        countervals = Counter(map(tuple,temp))
        for item in countervals:
            if len(item) == countervals[item]:
                # print "\n"+feature+" "+str(index),item,countervals[item]
                for ii in range(len(temp)):
                    if temp[ii] != list(item):
                        # print row[ii],col[ii],temp[ii],list(item)
                        for it in list(item):
                            if it in temp[ii]: elem_poss[str(row[ii])+','+str(col[ii])].remove(it)
    return elem_poss

def check_result(sudoku):
    row = [0]*9
    col = [0]*9
    sqr = [0]*9
    for ii in range(9):
        for jj in range(9):
            row[ii] += sudoku[ii][jj]
            col[jj] += sudoku[ii][jj]
            sqr[(ii/3)*3+jj/3] += sudoku[ii][jj]
    if row == [45]*9 and col == [45]*9 and sqr == [45]*9: return True
    else: return False

def solve_sudoku(sudoku):
    a_now = check_unfilled_cell(sudoku)
    a_prev = 0
    while a_prev != a_now:
        a_prev = a_now
        sudoku,elem_poss = check_cell_to_number(sudoku)
        sudoku = check_number_to_row_col_sq(elem_poss,sudoku)
        elem_poss = check_same('row',elem_poss)
        elem_poss = check_same('col',elem_poss)
        elem_poss = check_same('sqr',elem_poss)
        sudoku = check_number_to_row_col_sq(elem_poss,sudoku)
        a_now = check_unfilled_cell(sudoku)
    return sudoku,check_result(sudoku)

# result,status = solve_sudoku(sudoku)
# pprint.pprint(result)

from pandas.tools.plotting import scatter_matrix
from sklearn import cross_validation
from sklearn.metrics import classification_report, confusion_matrix, accuracy_score
from sklearn.linear_model import LogisticRegression, LinearRegression
from sklearn.tree import DecisionTreeClassifier
from sklearn.neighbors import KNeighborsClassifier
from sklearn.discriminant_analysis import LinearDiscriminantAnalysis
from sklearn.naive_bayes import GaussianNB
from sklearn.svm import SVC, LinearSVC
from sklearn.decomposition import PCA
from sklearn import datasets
from sklearn.cluster import KMeans, DBSCAN, AgglomerativeClustering, MeanShift
from sklearn.ensemble import RandomForestClassifier, BaggingClassifier, AdaBoostClassifier, GradientBoostingClassifier
from sklearn.neural_network import MLPClassifier
from scipy.cluster import hierarchy
from mpl_toolkits.mplot3d import Axes3D

# iris = datasets.load_iris()
# print iris.keys()
# class_label = iris.target_names
# feature_label = iris.feature_names
# feature_data = iris.data
# class_data = iris.target
# samples, features = feature_data.shape
# data = pd.DataFrame(feature_data, columns=feature_label)
# dataset = pd.DataFrame(feature_data, columns=feature_label)
# dataset['class'] = pd.Series(class_data, index=data.index)
# print data.describe()
# print data.groupby('class').size()
# data.plot(kind='box', subplots=True, layout=(2,2), sharex=False, sharey=False)
# data.hist()
# scatter_matrix(data)
# plt.show()

# regr = LinearRegression()
# regr.fit(X_train, Y_train)
# print regr.coef_

# iris = datasets.load_iris()
# X, Y = iris.data, iris.target
# fig = plt.figure(1, figsize=(8, 6))
# ax = Axes3D(fig, elev=-150, azim=110)
# X_reduced = PCA(n_components=3).fit_transform(iris.data)
# ax.scatter(X_reduced[:, 0], X_reduced[:, 1], X_reduced[:, 2], c=Y, cmap=plt.cm.Paired)
# ax.set_title("First three PCA directions")
# ax.set_xlabel("1st eigenvector")
# ax.w_xaxis.set_ticklabels([])
# ax.set_ylabel("2nd eigenvector")
# ax.w_yaxis.set_ticklabels([])
# ax.set_zlabel("3rd eigenvector")
# ax.w_zaxis.set_ticklabels([])
# plt.show()

def K_FOLD_CROSS_VALIDATION():
	iris = datasets.load_iris()
	X, Y = iris.data, iris.target
	X_train, X_test, Y_train, Y_test = cross_validation.train_test_split(X, Y, test_size=0.2)
	# kfold = cross_validation.KFold(n=len(X_train), n_folds=10, random_state=7)
	# cv_results = cross_validation.cross_val_score(model, X_train, Y_train, cv=kfold, scoring='accuracy')
	# mean = cv_results.mean()
	# std = cv_results.std()
	return X, X_train, X_test, Y, Y_train, Y_test

def GAUSSIAN_NAIVE_BAYES(X_train, X_test, Y_train, Y_test):
	gnb = GaussianNB()
	y_pred = gnb.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	return accuracy

def LINEAR_REGRESSION(X_train, X_test, Y_train, Y_test):
	lin = LinearRegression()
	y_pred_temp = lin.fit(X_train, Y_train).predict(X_test)
	# print lin.fit(X_train, Y_train).coef_
	y_pred = [1 for index,item in enumerate(y_pred_temp) if Y_test[index] != round(abs(item))]
	accuracy = sum(y_pred)
	return accuracy

def LOGISTIC_REGRESSION(X_train, X_test, Y_train, Y_test):
	lor = LogisticRegression()
	y_pred = lor.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	return accuracy

def LINEAR_DISCRIMINANT_ANALYSIS(X_train, X_test, Y_train, Y_test):
	lda = LinearDiscriminantAnalysis()
	y_pred = lda.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	return accuracy

def K_NEAREST_NEIGHBOURS(X_train, X_test, Y_train, Y_test):
	knn = KNeighborsClassifier()
	y_pred = knn.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	return accuracy

def DECISION_TREE(X_train, X_test, Y_train, Y_test):
	det = DecisionTreeClassifier()
	y_pred = det.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	return accuracy

def SUPPORT_VECTOR_MACHINE(X_train, X_test, Y_train, Y_test):
	svm = SVC()
	y_pred = svm.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	return accuracy

def PRINCIPAL_COMPONENT_ANALYSIS(X_train, X_test, Y_train, Y_test):
	pca = PCA(n_components=2)
	y_trans = pca.fit(X_train, Y_train)
	y_trans = pca.fit_transform(X_train, Y_train)
	return y_trans

def K_MEANS(X, clusters):
	kmeans = KMeans(n_clusters=clusters, random_state=0).fit(X)
	cluster_labels, cluster_centers = kmeans.labels_, kmeans.cluster_centers_
	# print kmeans.predict([[0, 0, 0, 0], [4, 4, 4, 4]])
	return cluster_labels, cluster_centers

def AGGLOMERATIVE(X, clusters, methods, metrics):
	agglo = AgglomerativeClustering(n_clusters=clusters, linkage=methods, affinity=metrics).fit_predict(X)
	return agglo

def HIERARCHIAL(X, maxdistance, maxclusters, methods, metrics, criterions):
	dist, clust_details = [], {}
	for i in range(len(X)-1):
	    for j in range(i+1,len(X)): dist.append(np.linalg.norm(X[i]-X[j]))
	clusters = hierarchy.linkage(dist, method=methods, metric=metrics)
	if criterions == 'distance': param = maxdistance
	if criterions == 'maxclust': param = maxclusters
	T = hierarchy.fcluster(clusters, param, criterion=criterions)
	for item in range(max(T)): clust_details[item+1] = []
	for index,item in enumerate(X): clust_details[T[index]].append(item)
	P = hierarchy.dendrogram(clusters)
	# plt.show()
	return T

def DBSCANS(X, maxeps, minsamples):
	dbscan = DBSCAN(eps=maxeps, min_samples = minsamples).fit(X)
	labels = dbscan.labels_
	return labels

def MEAN_SHIFT(X):
	mns = MeanShift().fit(X)
	cluster_labels, cluster_centers = mns.labels_, mns.cluster_centers_
	# print kmeans.predict([[0, 0, 0, 0], [4, 4, 4, 4]])
	return cluster_labels, cluster_centers

def ADA_BOOST(X_train, X_test, Y_train, Y_test, model):
	ada = AdaBoostClassifier(n_estimators=10, base_estimator=model, learning_rate=1)
	y_pred = ada.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	return accuracy

def GRADIENT_BOOST(X_train, X_test, Y_train, Y_test):
	ada = GradientBoostingClassifier(n_estimators=10, learning_rate=1)
	y_pred = ada.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	return accuracy

def BAGGING(X_train, X_test, Y_train, Y_test, model):
	bag = BaggingClassifier(base_estimator=model, n_estimators=10, bootstrap=True)
	y_pred = bag.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	return accuracy

def RANDOM_FOREST(X_train, X_test, Y_train, Y_test):
	rfr = RandomForestClassifier(n_estimators=10, criterion='gini', bootstrap=True)
	y_pred = rfr.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	# print rfr.fit(X_train, Y_train).feature_importances_
	return accuracy

def NEURAL_NETWORK(X_train, X_test, Y_train, Y_test):
	nn = MLPClassifier(solver='lbfgs', alpha=1e-5, hidden_layer_sizes=(15, 12))
	y_pred = nn.fit(X_train, Y_train).predict(X_test)
	accuracy = (Y_test != y_pred).sum()
	return accuracy

def ANALYSIS(Y_test, y_pred):
	print confusion_matrix(Y_test, y_pred)
	print accuracy_score(Y_test, y_pred)
	print classification_report(Y_test, y_pred)

# X, X_train, X_test, Y, Y_train, Y_test = K_FOLD_CROSS_VALIDATION()

# print GAUSSIAN_NAIVE_BAYES(X_train, X_test, Y_train, Y_test)
# print LINEAR_REGRESSION(X_train, X_test, Y_train, Y_test)
# print LOGISTIC_REGRESSION(X_train, X_test, Y_train, Y_test)
# print LINEAR_DISCRIMINANT_ANALYSIS(X_train, X_test, Y_train, Y_test)
# print K_NEAREST_NEIGHBOURS(X_train, X_test, Y_train, Y_test)
# print DECISION_TREE(X_train, X_test, Y_train, Y_test)
# print SUPPORT_VECTOR_MACHINE(X_train, X_test, Y_train, Y_test)
# print K_MEANS(X, 3)
# print AGGLOMERATIVE(X, 3, 'average', 'euclidean')
# print HIERARCHIAL(X, 2, 3, 'weighted', 'euclidean', 'distance')
# print DBSCANS(X, 0.5, 5)
# print MEAN_SHIFT(X)
# print ADA_BOOST(X_train, X_test, Y_train, Y_test, DecisionTreeClassifier())
# print GRADIENT_BOOST(X_train, X_test, Y_train, Y_test)
# print BAGGING(X_train, X_test, Y_train, Y_test, LinearDiscriminantAnalysis())
# print RANDOM_FOREST(X_train, X_test, Y_train, Y_test)
# print NEURAL_NETWORK(X_train, X_test, Y_train, Y_test)


# titanic_df = pd.read_csv("train.csv")
# test_df    = pd.read_csv("test.csv")
# print titanic_df.head()
# print titanic_df.info()
# print test_df.info()
# titanic_df = titanic_df.drop(['PassengerId','Name','Ticket'], axis=1)
# test_df    = test_df.drop(['Name','Ticket'], axis=1)
# titanic_df["Embarked"] = titanic_df["Embarked"].fillna("S")
# sns.factorplot('Embarked','Survived', data=titanic_df,size=4,aspect=3)
# fig, (axis1,axis2,axis3) = plt.subplots(1,3,figsize=(15,5))
# sns.countplot(x='Embarked', data=titanic_df, ax=axis1)
# sns.countplot(x='Survived', hue="Embarked", data=titanic_df, order=[1,0], ax=axis2)
# embark_perc = titanic_df[["Embarked", "Survived"]].groupby(['Embarked'],as_index=False).mean()
# sns.barplot(x='Embarked', y='Survived', data=embark_perc,order=['S','C','Q'],ax=axis3)
# embark_dummies_titanic  = pd.get_dummies(titanic_df['Embarked'])
# embark_dummies_titanic.drop(['S'], axis=1, inplace=True)
# embark_dummies_test  = pd.get_dummies(test_df['Embarked'])
# embark_dummies_test.drop(['S'], axis=1, inplace=True)
# titanic_df = titanic_df.join(embark_dummies_titanic)
# test_df    = test_df.join(embark_dummies_test)
# titanic_df.drop(['Embarked'], axis=1,inplace=True)
# test_df.drop(['Embarked'], axis=1,inplace=True)
# test_df["Fare"].fillna(test_df["Fare"].median(), inplace=True)
# titanic_df['Fare'] = titanic_df['Fare'].astype(int)
# test_df['Fare']    = test_df['Fare'].astype(int)
# fare_not_survived = titanic_df["Fare"][titanic_df["Survived"] == 0]
# fare_survived     = titanic_df["Fare"][titanic_df["Survived"] == 1]
# avgerage_fare = DataFrame([fare_not_survived.mean(), fare_survived.mean()])
# std_fare      = DataFrame([fare_not_survived.std(), fare_survived.std()])
# titanic_df['Fare'].plot(kind='hist', figsize=(15,3),bins=100, xlim=(0,50))
# avgerage_fare.index.names = std_fare.index.names = ["Survived"]
# avgerage_fare.plot(yerr=std_fare,kind='bar',legend=False)
# fig, (axis1,axis2) = plt.subplots(1,2,figsize=(15,4))
# axis1.set_title('Original Age values - Titanic')
# axis2.set_title('New Age values - Titanic')
# average_age_titanic   = titanic_df["Age"].mean()
# std_age_titanic       = titanic_df["Age"].std()
# count_nan_age_titanic = titanic_df["Age"].isnull().sum()
# average_age_test   = test_df["Age"].mean()
# std_age_test       = test_df["Age"].std()
# count_nan_age_test = test_df["Age"].isnull().sum()
# rand_1 = np.random.randint(average_age_titanic - std_age_titanic, average_age_titanic + std_age_titanic, size = count_nan_age_titanic)
# rand_2 = np.random.randint(average_age_test - std_age_test, average_age_test + std_age_test, size = count_nan_age_test)
# titanic_df['Age'].dropna().astype(int).hist(bins=70, ax=axis1)
# titanic_df["Age"][np.isnan(titanic_df["Age"])] = rand_1
# test_df["Age"][np.isnan(test_df["Age"])] = rand_2
# titanic_df['Age'] = titanic_df['Age'].astype(int)
# test_df['Age']    = test_df['Age'].astype(int)
# titanic_df['Age'].hist(bins=70, ax=axis2)
# facet = sns.FacetGrid(titanic_df, hue="Survived",aspect=4)
# facet.map(sns.kdeplot,'Age',shade= True)
# facet.set(xlim=(0, titanic_df['Age'].max()))
# facet.add_legend()
# fig, axis1 = plt.subplots(1,1,figsize=(18,4))
# average_age = titanic_df[["Age", "Survived"]].groupby(['Age'],as_index=False).mean()
# sns.barplot(x='Age', y='Survived', data=average_age)
# titanic_df.drop("Cabin",axis=1,inplace=True)
# test_df.drop("Cabin",axis=1,inplace=True)
# titanic_df['Family'] =  titanic_df["Parch"] + titanic_df["SibSp"]
# titanic_df['Family'].loc[titanic_df['Family'] > 0] = 1
# titanic_df['Family'].loc[titanic_df['Family'] == 0] = 0
# test_df['Family'] =  test_df["Parch"] + test_df["SibSp"]
# test_df['Family'].loc[test_df['Family'] > 0] = 1
# test_df['Family'].loc[test_df['Family'] == 0] = 0
# titanic_df = titanic_df.drop(['SibSp','Parch'], axis=1)
# test_df    = test_df.drop(['SibSp','Parch'], axis=1)
# fig, (axis1,axis2) = plt.subplots(1,2,sharex=True,figsize=(10,5))
# sns.countplot(x='Family', data=titanic_df, order=[1,0], ax=axis1)
# family_perc = titanic_df[["Family", "Survived"]].groupby(['Family'],as_index=False).mean()
# sns.barplot(x='Family', y='Survived', data=family_perc, order=[1,0], ax=axis2)
# axis1.set_xticklabels(["With Family","Alone"], rotation=0)

# def get_person(passenger):
#     age,sex = passenger
#     return 'child' if age < 16 else sex

# titanic_df['Person'] = titanic_df[['Age','Sex']].apply(get_person,axis=1)
# test_df['Person']    = test_df[['Age','Sex']].apply(get_person,axis=1)
# titanic_df.drop(['Sex'],axis=1,inplace=True)
# test_df.drop(['Sex'],axis=1,inplace=True)
# person_dummies_titanic  = pd.get_dummies(titanic_df['Person'])
# person_dummies_titanic.columns = ['Child','Female','Male']
# person_dummies_titanic.drop(['Male'], axis=1, inplace=True)
# person_dummies_test  = pd.get_dummies(test_df['Person'])
# person_dummies_test.columns = ['Child','Female','Male']
# person_dummies_test.drop(['Male'], axis=1, inplace=True)
# titanic_df = titanic_df.join(person_dummies_titanic)
# test_df    = test_df.join(person_dummies_test)
# fig, (axis1,axis2) = plt.subplots(1,2,figsize=(10,5))
# sns.countplot(x='Person', data=titanic_df, ax=axis1)
# person_perc = titanic_df[["Person", "Survived"]].groupby(['Person'],as_index=False).mean()
# sns.barplot(x='Person', y='Survived', data=person_perc, ax=axis2, order=['male','female','child'])
# titanic_df.drop(['Person'],axis=1,inplace=True)
# test_df.drop(['Person'],axis=1,inplace=True)
# sns.factorplot('Pclass','Survived',order=[1,2,3], data=titanic_df,size=5)
# pclass_dummies_titanic  = pd.get_dummies(titanic_df['Pclass'])
# pclass_dummies_titanic.columns = ['Class_1','Class_2','Class_3']
# pclass_dummies_titanic.drop(['Class_3'], axis=1, inplace=True)
# pclass_dummies_test  = pd.get_dummies(test_df['Pclass'])
# pclass_dummies_test.columns = ['Class_1','Class_2','Class_3']
# pclass_dummies_test.drop(['Class_3'], axis=1, inplace=True)
# titanic_df.drop(['Pclass'],axis=1,inplace=True)
# test_df.drop(['Pclass'],axis=1,inplace=True)
# titanic_df = titanic_df.join(pclass_dummies_titanic)
# test_df    = test_df.join(pclass_dummies_test)
# X_train = titanic_df.drop("Survived",axis=1)
# Y_train = titanic_df["Survived"]
# X_test  = test_df.drop("PassengerId",axis=1).copy()
# logreg = LogisticRegression()
# logreg.fit(X_train, Y_train)
# Y_pred = logreg.predict(X_test)
# logreg.score(X_train, Y_train)
# random_forest = RandomForestClassifier(n_estimators=100)
# random_forest.fit(X_train, Y_train)
# Y_pred = random_forest.predict(X_test)
# random_forest.score(X_train, Y_train)
# coeff_df = DataFrame(titanic_df.columns.delete(0))
# coeff_df.columns = ['Features']
# coeff_df["Coefficient Estimate"] = pd.Series(logreg.coef_[0])
# print coeff_df
# plt.show()

# from sklearn.datasets import fetch_20newsgroups
# from sklearn.feature_extraction.text import CountVectorizer
# from sklearn.feature_extraction.text import TfidfTransformer
# from sklearn.naive_bayes import MultinomialNB
# from sklearn.pipeline import Pipeline
# from sklearn.linear_model import SGDClassifier
# from sklearn import metrics

# twenty_train = fetch_20newsgroups(subset='train', shuffle=True, random_state=42)
# train_data = twenty_train.data
# train_class = twenty_train.target
# class_names = twenty_train.target_names
# twenty_test = fetch_20newsgroups(subset='test', shuffle=True, random_state=42)
# test_data = twenty_test.data
# test_class = twenty_test.target

# count_vect = CountVectorizer()
# X_train_counts = count_vect.fit_transform(train_data)

# tf_transformer = TfidfTransformer(use_idf=False).fit(X_train_counts)
# X_train_tf = tf_transformer.transform(X_train_counts)

# tfidf_transformer = TfidfTransformer()
# X_train_tfidf = tfidf_transformer.fit_transform(X_train_counts)

# clf = MultinomialNB().fit(X_train_tfidf, train_class)

# td = ['aa bb cc dd', 'aa bb cc dd aa bb cc dd', 'aa bb cc dd ee ff', 'bb cc']
# count_vect = CountVectorizer()
# tt = count_vect.fit_transform(td).toarray()
# # count_vect.vocabulary_

# for index in range(X_train_counts.shape[0]):
#     # if index%1000 == 0: print('>>>>>',index)
#     val = sum(X_train_counts[index,:].toarray()[0])

# docs_new = ['God is love', 'OpenGL on the GPU is fast']
# X_new_counts = count_vect.transform(docs_new)
# X_new_tfidf = tfidf_transformer.transform(X_new_counts)
# predicted = clf.predict(X_new_tfidf)

# for doc, category in zip(docs_new, predicted): print('%r => %s' % (doc, twenty_train.target_names[category]))

# text_clf = Pipeline([('vect', CountVectorizer()), ('tfidf', TfidfTransformer()), ('clf', MultinomialNB())])
# text_clf = text_clf.fit(twenty_train.data, twenty_train.target)
# predicted = text_clf.predict(test_data)
# acc = round((np.mean(predicted == test_class)*100))
# print acc

# text_clf = Pipeline([('vect', CountVectorizer()), ('tfidf', TfidfTransformer()), ('clf', SGDClassifier(loss='hinge', penalty='l2',alpha=1e-3, n_iter=5, random_state=42))])
# text_clf = text_clf.fit(twenty_train.data, twenty_train.target)
# predicted = text_clf.predict(test_data)
# acc = round((np.mean(predicted == test_class)*100))
# print acc

# import PyPDF2,os,csv

# filenames = os.listdir("Walmart RP/Classified")
# # write_filename = "test.csv"
# # writer = csv.writer(open(write_filename, 'w'),delimiter=';')
# # writer.writerow(['claimNum','date','to','suppNum','claimType','claimExp','AmtPend','mgrName'])
# counts = 0
# for filename in filenames: 
#     if filename.endswith('.pdf'):
#         counts += 1
#         pdfFileObj = open('Walmart RP/Classified/'+filename, 'rb')
#         pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
#         numPage = pdfReader.numPages
#         pageObj = pdfReader.getPage(0)
#         pageText = pageObj.extractText().split('\n')
#         date = pageText[16]
#         suppName = pageText[17]
#         claimNum = pageText[18]
#         suppNum = pageText[19]
#         claimType = pageText[21]
#         claimExp = pageText[23]
#         AmtPend = pageText[24].replace(',','')
#         mgrName = pageText[40]
#         # for index,item in enumerate(pageText): print index,item
#         print filename
#         # writer.writerow([claimNum,date,to,suppNum,claimType,claimExp,AmtPend,mgrName])
#         # pageObj = pdfReader.getPage(1)
#         # pageText = pageObj.extractText().split('\n')
#         # print counts,len(pageText),claimType

#         # pageObj = pdfReader.getPage(4)
#         # pageText = pageObj.extractText()

#     if counts == 1: break
# # writer = csv.writer(open('write_filename', 'w'),delimiter=';')

# from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
# from pdfminer.converter import TextConverter
# from pdfminer.layout import LAParams
# from pdfminer.pdfpage import PDFPage
# from cStringIO import StringIO

# def convert_pdf_to_txt(path):
#     rsrcmgr = PDFResourceManager()
#     retstr = StringIO()
#     codec = 'utf-8'
#     laparams = LAParams()
#     device = TextConverter(rsrcmgr, retstr, codec=codec, laparams=laparams)
#     fp = file(path, 'rb')
#     interpreter = PDFPageInterpreter(rsrcmgr, device)
#     password = ""
#     maxpages = 0
#     caching = True
#     pagenos=set()

#     for pagenum,page in enumerate(PDFPage.get_pages(fp, pagenos, maxpages=maxpages, password=password,caching=caching, check_extractable=True)):
#         print pagenum,
#         if pagenum > 2: interpreter.process_page(page)
#         # interpreter.process_page(page)
#     print ''
#     text = retstr.getvalue()

#     fp.close()
#     device.close()
#     retstr.close()
#     return text

# catchword0 = ['corrections','errors','price adjustments','price differences','current cost','price protection','new cost','new retail cost','move forward cost']
# catchword1 = ['Price Change', 'Effective Date', 'Invoice Error', 'Pricing Error']
# catchword2 = ['Promotional Off', 'Mark Down', 'Rebate', 'COOP']
# catchword3 = ['Cost decline', 'Retail Price', 'Inventory', 'Effective Date']

# filelist = ['415620112.pdf','415621899.pdf','415621918.pdf','415621922.pdf','415621930.pdf','415621935.pdf','415622002.pdf','415622061.pdf','415622070.pdf','415622071.pdf','415622084.pdf','415622096.pdf','415622109.pdf','415622204.pdf']
# filelist = ['415620014.pdf']
# filenames = os.listdir("Walmart RP/Classified")
# for filename in filenames:
#     if filename in filelist:
#     # if filename.endswith('.pdf'):
#         text = convert_pdf_to_txt('Walmart RP/Classified/'+filename).replace('\xc2\xa0',' ').replace('\xe2\x80\x9d','').replace('\xe2\x80\x93','-').replace('\xc2\xad','-').replace('\xe2\x80\x99','').replace('\x0c','\n').replace('\xe2\x80\x98','').replace('\xe2\x80\xa6','...').lower()
#         # catch0 = [word.lower() for word in catchword0 if (word.lower() in text)]
#         # catch1 = [word.lower() for word in catchword1 if (word.lower() in text)]
#         # catch2 = [word.lower() for word in catchword2 if (word.lower() in text)]
#         # catch3 = [word.lower() for word in catchword3 if (word.lower() in text)]
#         # print catch0
#         # print catch1
#         # print catch2
#         # print catch3
#         # text = " ".join(text.split())
#         text = text.split('\n')
#         print text
#         # break

from googlefinance import getQuotes
import subprocess,fbchat,time,mp3play,editdistance,os,sys,json,itertools,feedparser,googlemaps,pyautogui,urllib,urllib2,wikipedia,pytz
from weather import Weather
from pycricbuzz import Cricbuzz
from datetime import datetime
from periodic import element,table
from geopy import geocoders

def get_file(msg,folder_name,ext):
    filenames = os.listdir(folder_name)
    closest_name = ''
    mindist = 99
    for filename in filenames:
        filename = filename.lower()
        if filename.endswith(ext):
            dist = editdistance.eval(msg,filename.replace(ext,''))
            if dist < mindist: 
                closest_name = filename
                mindist = dist
    return closest_name    

def get_all_file(folder_name,ext):
    filenames = os.listdir(folder_name)
    file_names = []
    for filename in filenames:
        if filename.lower().endswith(ext):
            file_names.append(folder_name+filename)
    return file_names

def get_share(code):
    data = getQuotes('BOM:'+str(code))
    price = data[0]["LastTradePrice"]
    return 'Latest price of '+code+': '+str(price)

def music_player(last_msg,msg,song_name,clip,folder_name):
    if last_msg == msg: return last_msg,song_name,clip
    else: last_msg = msg
    if msg not in ['PLAY','PAUSE','UNPAUSE','STOP']:
        if 'VOLUME' in msg:
            if song_name == '': return last_msg,song_name,clip
            vol = int(msg.replace('VOLUME ',''))
            if vol > -1 and vol < 101: clip.volume(vol)
        else:
            song_name = get_file(msg.lower(),folder_name,'.mp3')
            print msg,song_name
            clip = mp3play.load(folder_name+song_name)
    else:
        if song_name == '': return last_msg,song_name,clip
        if msg == 'PLAY': clip.play()
        if msg == 'PAUSE': clip.pause()
        if msg == 'UNPAUSE': clip.unpause()
        if msg == 'STOP': clip.stop()
    return last_msg,song_name,clip

def get_last_msg(friend_id):
    last_messages = client.getThreadInfo(friend_id,100)
    for index,item in enumerate(last_messages):
        if friend_id == int(item.author.replace('fbid:','')): return last_messages[index]

def solving_equation(numbs):
    numb = map(int,numbs.split(','))
    numb_perm = list(itertools.permutations(numb[0:-1],4))
    sign_perm = list(itertools.product(['+','-','*','/'],repeat=3))
    for a1 in range(24):
        for b1 in range(64):
            string1 = str(float(numb_perm[a1][0]))+sign_perm[b1][0]+str(float(numb_perm[a1][1]))+sign_perm[b1][1]+str(float(numb_perm[a1][2]))+sign_perm[b1][2]+str(float(numb_perm[a1][3]))
            if eval(string1) == numb[-1]: return string1.replace('.0','')
            string2 = '((('+str(float(numb_perm[a1][0]))+sign_perm[b1][0]+str(float(numb_perm[a1][1]))+')'+sign_perm[b1][1]+str(float(numb_perm[a1][2]))+')'+sign_perm[b1][2]+str(float(numb_perm[a1][3]))+')'
            if eval(string2) == numb[-1]: return string2.replace('.0','')

def conv_f_to_c(val): return str(round((int(val)-32)*5.0/9.0,2))

def get_weather(city):
    weather = Weather()
    data = weather.lookup_by_location(city).print_obj()

    astronomy = data['astronomy']
    atmosphere = data['atmosphere']
    wind = data['wind']
    units = data['units']

    item = data['item']
    condition = item['condition']
    forecast = item['forecast']
    title = item['title']
    latitude = item['lat']
    longitude = item['long']

    result = title+'\nTemperature: '+conv_f_to_c(condition['temp'])+'*C'
    result += '\nWeather: '+condition['text']
    result += '\nHumidity: '+atmosphere['humidity']+'%'
    result += '\nPressure: '+atmosphere['pressure']+'millibars'
    result += '\nVisibility :'+atmosphere['visibility']+'Km'
    result += '\nSunrise: '+astronomy['sunrise']
    result += '\nSunset: '+astronomy['sunset']
    result += '\nWind: '+wind['speed']+'Km/hr'
    result += '\nDirection: '+wind['direction']+'*'
    result += '\nChill: '+conv_f_to_c(wind['chill'])+'*C'
    result += '\n\nFORECAST'
    
    for item in forecast: result += '\n'+item['date']+' '+item['day']+' '+conv_f_to_c(item['high'])+' '+conv_f_to_c(item['low'])+' '+item['text']
    return result

def get_cricket(msg):
    c = Cricbuzz()
    matches = c.matches()
    result = ''
    if msg == 'ALL':
        for match in matches:
            mid = match['id']
            desc = match['mchdesc']
            state = match['mchstate']
            source = match['srs']
            status = match['status']
            kind = match['type']        
            # result += mid+' '+desc+' '+source+'\n'
            result += desc+' ('+mid+')\n'
        return result
    else:
        for match in matches:
            mid = match['id']
            if mid != msg: continue
            else: 
                desc = match['mchdesc']
                state = match['mchstate']
                source = match['srs']
                status = match['status']
                kind = match['type']
                result += desc+' '+source
                result += '\nSTATUS : '+status
                if state == 'preview': return result+" (+05:30)"
                temp = c.scorecard(mid)
                scorecard = temp['scorecard']
                team1 = temp['squad'][0]
                team2 = temp['squad'][1]
                temp = c.livescore(mid)
                job = ['batting','bowling','batting','bowling']
                for index,scores in enumerate(scorecard):
                    short = temp[job[index]]['score'][0]
                    result +=  '\n\nINNINGS '+str(len(scorecard)-index)+' '+scores['batteam']+' '+short['runs']+'/'+short['wickets']+'('+short['overs']+')\n'
                    for batsman in scores['batcard']: result += '\n'+batsman['name']+' '+batsman['dismissal']+' '+batsman['runs']+' '+batsman['balls']+' '+batsman['fours']+' '+batsman['six']
                    result += '\n'
                    for bowler in scores['bowlcard']: result += '\n'+bowler['name']+' '+bowler['overs']+' '+bowler['maidens']+' '+bowler['runs']+' '+bowler['wickets']
                return result

def get_route(msg):
    msg_parts = msg.split(",")
    restype = msg_parts[0]
    source = msg_parts[1]
    destination = msg_parts[2]
    if msg_parts[3] == "NOW": clock = time.time()
    else: clock = int(time.mktime(datetime.strptime(str(msg_parts[3]), "%d/%m/%y %H:%M").timetuple()))
    key = 'AIzaSyB0vDdauLIixfUqauZIblqxbzP3hFHLCMw'
    mode = 'transit'
    unit = 'metric'
    gmaps = googlemaps.Client(key=key)
    directions_result = gmaps.directions(source,destination,mode=mode,departure_time=clock,alternatives=True)
    total_ways = len(directions_result)
    result_long, result_short = source+' TO '+destination+'\n\n', source+' TO '+destination+'\n\n'
    for index,item in enumerate(directions_result):
        voice = "OPTION "+str(index+1)
        result_short += '\n'+"OPTION "+str(index+1)+'\n'
        result_long += '\n'+"OPTION "+str(index+1)+'\n'
        route = item['legs'][0]
        for item in route['steps']:
            instruction = item['html_instructions']
            distance = item['distance']['text']
            duration = item['duration']['text']
            mode = item['travel_mode']
            voice = mode+' '+instruction+' '+distance+' '+duration.replace('mins','minutes')
            result_short += mode+' '+instruction+' '+distance+' '+duration.replace('mins','minutes')+'\n'
            result_long += mode+' '+instruction+' '+distance+' '+duration.replace('mins','minutes')+'\n'
            if mode == 'WALKING':
                text = item['steps']
                for item in text:
                    walk_distance = item['distance']['text']
                    walk_duration = item['duration']['text']                
                    try:
                        walk_instruction = item['html_instructions'].replace('<b>','').replace('</b>','').replace('</div>','').replace('&nbsp;',' ')
                        walk_instruction = walk_instruction.replace('<div style="font-size:0.9em">',' , ')
                        walk_instruction = walk_instruction+' '+walk_distance+' '+walk_duration.replace('mins','minutes')
                        result_long += '->'+walk_instruction+'\n'
                    except: pass
                
            if mode == 'TRANSIT':
                text = item['transit_details']
                num_stops = text['num_stops']
                arrival_stop = text['arrival_stop']['name']
                arrival_time = text['arrival_time']['text']
                departure_stop = text['departure_stop']['name']
                departure_time = text['departure_time']['text']
                vehicle_destination = text['headsign']
                vehicle_route = text['line']['name']
                try: vehicle_number = text['line']['short_name']
                except: vehicle_number = ''
                vehicle_type = text['line']['vehicle']['name']
                result_long += '->'+vehicle_type+' '+vehicle_number+' '+vehicle_route+' '+vehicle_destination+' '+departure_stop+' '+departure_time+' '+arrival_stop+' '+arrival_time+' '+str(num_stops)+'\n'

    if restype == 'LONG': return result_long
    else: return result_short

def get_news(msg):
    news_dict = {'TOP':'ndtvnews-top-stories','LATEST':'ndtvnews-latest','TRENDING':'ndtvnews-trending-news','INDIA':'ndtvnews-india-news','WORLD':'ndtvnews-world-news','BUSINESS':'ndtvprofit-latest','MOVIES':'ndtvmovies-latest','CRICKET':'ndtvsports-cricket','SPORTS':'ndtvsports-latest','TECH':'gadgets360-latest'}
    result = msg+' NEWS\n\n'
    data = feedparser.parse('http://feeds.feedburner.com/'+news_dict[msg])
    for index,item in enumerate(data['entries']):
        title = item['title']
        link = item['link']
        result += str(index+1)+'. '+link+'\n\n'
    return result   

def get_sports(msg):
    sports_dict = {'ALL':'all','CRICKET':'cricket','FOOTBALL':'football','TENNIS':'tennis'}
    result = msg+' NEWS\n\n'
    data = feedparser.parse('https://sports.ndtv.com/rss/'+sports_dict[msg])
    for index,item in enumerate(data['entries'][0:20]):
        title = item['title']
        link = item['link']
        result += str(index+1)+'. '+link+'\n\n'
    return result

def vlc_player(msg,folder_name,ext):
    if msg.split(" ")[0] == 'PLAY':
        subprocess.Popen("taskkill /IM vlc.exe")
        filenames = msg.replace("PLAY ","")
        if filenames == 'ALL': file_names = get_all_file(folder_name,ext)
        else: file_names = [folder_name+item+ext for item in msg.replace("PLAY ","").split(",")]
        subprocess.Popen([vlc_player_path]+file_names)
        return
    if len(msg.split(" ")) == 1:
        if msg == 'PAUSE': pyautogui.press('space')
        if msg == 'STOP': pyautogui.press('s')
        if msg == 'RESTART': pyautogui.press('p')
        if msg == 'QUIT': pyautogui.hotkey('ctrl', 'q')
        if msg == 'NEUTRAL': pyautogui.press('=')
        if msg == 'FASTEST': pyautogui.press('add')
        if msg == 'SLOWEST': pyautogui.press('subtract')
        if msg == 'MUTE': pyautogui.press('m')
        if msg == 'NEXT': pyautogui.press('n')
        if msg == 'PREVIOUS': pyautogui.press('p')
        if msg == 'FULLSCREEN': pyautogui.press('f')
        return
    if len(msg.split(" ")) == 2:
        msg_type = msg.split(" ")[0]
        msg_val = int(msg.split(" ")[1])
        if msg_type == 'FASTER':
            for ind in range(msg_val): pyautogui.press(']')
        if msg_type == 'SLOWER':
            for ind in range(msg_val): pyautogui.press('[')
        if msg_type == 'LOUDER':
            for ind in range(msg_val/5): pyautogui.hotkey('ctrl', 'up')
        if msg_type == 'QUIETER':
            for ind in range(msg_val/5): pyautogui.hotkey('ctrl', 'down')
        if msg_type == 'FORWARD':
            for ind in range(msg_val): pyautogui.hotkey('alt','right')
        if msg_type == 'BACKWARD':
            for ind in range(msg_val): pyautogui.hotkey('alt','left')
        if msg_type == 'FASTFORWARD':
            for ind in range(msg_val): pyautogui.hotkey('ctrl','right')
        if msg_type == 'FASTBACKWARD':
            for ind in range(msg_val): pyautogui.hotkey('ctrl','left')
        if msg_type == 'VOLUME':
            pyautogui.press('space')
            for ind in range(40): pyautogui.hotkey('ctrl', 'down')
            for ind in range(msg_val/5): pyautogui.hotkey('ctrl', 'up')
            pyautogui.press('space')
        return

def solving_sudoku(msg):
    sudoku = [list(map(int,list(msg[index*9:(index+1)*9]))) for index in range(9)]
    solution,status = solve_sudoku(sudoku)
    if status == False: return 'Unable to solve'
    result = ''
    for item in solution: result += str(item).replace('[','').replace(']','')+'\n'
    return result

def get_wiki(msg):
    try:
        temp = wikipedia.page(msg.split(',')[1])
        if msg.split(',')[0] == 'LONG': return temp.content
        else: return temp.summary
    except: return 'Unable to find'

def get_pin_map(msg):
    result = ''
    if msg.isdigit() == True:
        for item in requests.get("https://pincode.saratchandra.in/api/pincode/"+msg).json()['data']:
            result += 'PINCODE: '+str(item['pincode'])+' PO Name: '+str(item['office_name'])+' Region: '+str(item['region_name'])+' District: '+str(item['district'])+' State: '+str(item['state_name'])+'\n\n'
    else:
        for item in requests.get("https://pincode.saratchandra.in/api/pincode/",params={"city": msg.title()}).json()['data']:
            result += 'PINCODE: '+str(item['pincode'])+' PO Name: '+str(item['office_name'])+' Region: '+str(item['region_name'])+' District: '+str(item['district'])+' State: '+str(item['state_name'])+'\n\n'    
    return result

def get_element(msg):
    if msg == 'TABLE': return table
    else: 
        elem = element(msg)
        return 'Element: '+elem.name+' Symbol: '+elem.symbol+' Atomic no: '+str(elem.atomic)+' Atomic Mass: '+str(elem.mass)

def get_location_time(msg):
    g = geocoders.GoogleV3()
    place, (lat, lng) = g.geocode(msg)
    url = 'http://api.geonames.org/timezoneJSON?lat='+str(lat)+'&lng='+str(lng)+'&username=demo'
    resp = requests.get(url).json()
    clock = resp['time'].split(" ")
    result = 'Current time at '+place+'\nDate: '+clock[0]+' Time: '+clock[1]
    return result

def get_msg(current_time,friend_name,msgtype):
    hr = int(datetime.fromtimestamp(int(current_time/1000)).strftime('%H'))
    timeslot = {'Morning':[6,7,8,9,10,11],'Afternoon':[12,13,14,15],'Evening':[16,17,18],'Night':[0,1,2,3,4,5,19,20,21,22,23]}
    for item in timeslot.keys():
        if hr in timeslot[item]:
            start_msg = 'Good '+item+' '+friend_name.split(" ")[0]+'!!!\nWelcome to ChatBot'
            end_msg = 'Bye '+friend_name.split(" ")[0]+'!!!\nHave a nice '+item+' ahead'
    if msgtype == 'start': return start_msg
    else: return end_msg

def get_help(msg):
    result = ''
    if msg == 'HELP':
        result += 'Hello '+friend_name.split(" ")[0]+', how can I help you?\n\n'
        result += 'Type HELP <topic> for details\n\n'
        result += 'Ping, Music, Share, Solve, Weather, Cricket, Route, News, Video, Audio, Sudoku, Wiki, Sports, All, Pincode, Periodic, Time'
        return result

# music_folder_name = 'C:\Users/sourya.poddar/Downloads/Personal/BANGLA/MANNA DE/'
# video_folder_name = "C:\Users/sourya.poddar/Downloads/Personal/"
# vlc_player_path = "D:/vlc-2.2.4/vlc.exe"
# fb_id = 'souryapoddar290990@gmail.com'
# fb_pass = 'temp1234'
# friend_name = "Sourya Poddar"
# fb_pass = 'sourya1000'
# fb_id = 'sourya290990@yahoo.com'
# friend_name = "Debojyoti Dey"
# client = fbchat.Client(fb_id,fb_pass)
# friends = client.getUsers(friend_name)
# friend = friends[0]
# friend_id = friend.uid
# last_music_msg = ''
# song_name = ''
# current_time = time.time()*1000
# start_msg = get_msg(current_time,friend_name,'start')
# last_msg_id = ''
# clip = ''
# sent = client.send(friend_id,start_msg)

# no_reply = ['Sorry, I dont get you','I have no answer for that','Can you reframe your query?']
# instructions = ['HI','HELLO','BYE','HELP','PING','MUSIC','SHARE','SOLVE','WEATHER','CRICKET','ROUTE','NEWS','VIDEO','AUDIO','SUDOKU','WIKI','SPORTS','PINCODE','PERIODIC','TIME']
# while True:
#     # break
#     try:
#         send_msg = ''
#         last_messages = get_last_msg(friend_id)
#         msg = last_messages.body.upper()
#         tst = last_messages.timestamp
#         eml = last_messages.author
#         mid = last_messages.message_id
#         print msg,tst,mid
#         if current_time > tst or mid == last_msg_id: pass
#         else:
#             msg_type = msg.split(" ")[0]
#             if msg_type not in instructions: send_msg = no_reply[np.random.randint(0,3,1)[0]]+'\nFor assistance, type HELP'
#             if msg_type in ['HI','HELLO']: send_msg = 'Hi '+friend_name.split(" ")[0]
#             if msg_type in ['BYE','GOODBYE','TATA']: send_msg = get_msg(current_time,friend_name,'stop')
#             if msg_type == 'HELP': send_msg = get_help(msg)
#             if msg_type == 'PING': send_msg = 'PING TIME '+str(round(time.time()-tst/1000,2))
#             if msg_type == 'MUSIC': last_music_msg,song_name,clip = music_player(last_music_msg,msg.replace("MUSIC ",""),song_name,clip,music_folder_name)
#             if msg_type == 'SHARE': send_msg = get_share(msg.split(" ")[1])
#             if msg_type == 'SOLVE': send_msg = solving_equation(msg.split(" ")[1])
#             if msg_type == 'WEATHER': send_msg = get_weather(msg.replace("WEATHER ",""))
#             if msg_type == 'CRICKET': send_msg = get_cricket(msg.split(" ")[1])
#             if msg_type == 'ROUTE': send_msg = get_route(msg.replace("ROUTE ",""))
#             if msg_type == 'NEWS': send_msg = get_news(msg.replace("NEWS ","")) 
#             if msg_type == 'VIDEO': vlc_player(msg.replace("VIDEO ",""),video_folder_name,'.mp4')
#             if msg_type == 'AUDIO': vlc_player(msg.replace("AUDIO ",""),music_folder_name,'.mp3')
#             if msg_type == 'SUDOKU': send_msg = solving_sudoku(msg)
#             if msg_type == 'WIKI': send_msg = get_wiki(msg.replace("WIKI ",""))
#             if msg_type == 'SPORTS': send_msg = get_sports(msg.replace("SPORTS ",""))
#             if msg_type == 'PINCODE': send_msg = get_pin_map(msg.replace("PINCODE ",""))
#             if msg_type == 'PERIODIC': send_msg = get_element(msg.replace("PERIODIC ",""))
#             if msg_type == 'TIME': send_msg = get_location_time(msg.replace("TIME ",""))
#         if send_msg != '':
#             sent = client.send(friend_id,send_msg)
#             if sent: last_msg_id = mid
#         else: last_msg_id = mid
#     except Exception,e: 
#         print e
#         last_msg_id = mid
#     time.sleep(1)



# msg = 'hoyto tomari'
# song_name = get_music(msg,folder_name)
# process_name = '"C:/Program Files (x86)/Windows Media Player/wmplayer.exe"'
# file_name = '"'+folder_name+song_name+'"'
# subprocess_call = process_name+' '+file_name
# print subprocess_call
# process = subprocess.Popen(subprocess_call)
# time.sleep(10)
# process.terminate()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
# ########################################################################
# prs = Presentation()
# ########################################################################
# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]
# title.text = "Hello, World!"
# subtitle.text = "python-pptx was here!"
# ########################################################################
# bullet_slide_layout = prs.slide_layouts[1]
# slide = prs.slides.add_slide(bullet_slide_layout)
# shapes = slide.shapes
# title_shape = shapes.title
# body_shape = shapes.placeholders[1]
# title_shape.text = 'Adding a Bullet Slide'
# tf = body_shape.text_frame
# tf.text = 'Find the bullet slide layout'
# p = tf.add_paragraph()
# p.text = 'Use _TextFrame.text for first bullet'
# p.level = 1
# p = tf.add_paragraph()
# p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
# p.level = 2
# ########################################################################
# blank_slide_layout = prs.slide_layouts[6]
# slide = prs.slides.add_slide(blank_slide_layout)
# left = top = width = height = Inches(1)
# txBox = slide.shapes.add_textbox(left, top, width, height)
# tf = txBox.text_frame
# tf.text = "This is text inside a textbox"
# p = tf.add_paragraph()
# p.text = "This is a second paragraph that's bold"
# p.font.bold = True
# p = tf.add_paragraph()
# p.text = "This is a third paragraph that's big"
# p.font.size = Pt(40)
# ########################################################################
# img_path = 'test_img.png'
# blank_slide_layout = prs.slide_layouts[6]
# slide = prs.slides.add_slide(blank_slide_layout)
# left = top = Inches(1)
# pic = slide.shapes.add_picture(img_path, left, top)
# left = Inches(5)
# height = Inches(5.5)
# pic = slide.shapes.add_picture(img_path, left, top, height=height)
# ########################################################################
# title_only_slide_layout = prs.slide_layouts[5]
# slide = prs.slides.add_slide(title_only_slide_layout)
# shapes = slide.shapes
# shapes.title.text = 'Adding an AutoShape'
# left = Inches(0.93)
# top = Inches(3.0)
# width = Inches(1.75)
# height = Inches(1.0)
# shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
# shape.text = 'Step 1'
# left = left + width - Inches(0.4)
# width = Inches(2.0)
# for n in range(2, 6):
#     shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
#     shape.text = 'Step %d' % n
#     left = left + width - Inches(0.4)
# ########################################################################
# title_only_slide_layout = prs.slide_layouts[5]
# slide = prs.slides.add_slide(title_only_slide_layout)
# shapes = slide.shapes
# shapes.title.text = 'Adding a Table'
# rows = cols = 2
# left = top = Inches(2.0)
# width = Inches(6.0)
# height = Inches(0.8)
# table = shapes.add_table(rows, cols, left, top, width, height).table
# table.columns[0].width = Inches(2.0)
# table.columns[1].width = Inches(4.0)
# table.cell(0, 0).text = 'Foo'
# table.cell(0, 1).text = 'Bar'
# table.cell(1, 0).text = 'Baz'
# table.cell(1, 1).text = 'Qux'
# ########################################################################
# prs.save('test1.pptx')

# nltk.download()
# text = word_tokenize("And now for something completely different")
# text = word_tokenize("They refuse to permit us to obtain the refuse permit")
# print nltk.pos_tag(text)

# text = nltk.Text(word.lower() for word in nltk.corpus.brown.words())
# print text.similar('woman')
# print text.similar('bought')
# print text.similar('over')
# print text.similar('the')

# tagged_token = nltk.tag.str2tuple('fly/NN')
# print nltk.corpus.brown.tagged_words()
# print nltk.corpus.brown.tagged_words(tagset='universal')
# brown_news_tagged = nltk.corpus.brown.tagged_words(categories='news', tagset='universal')
# tag_fd = nltk.FreqDist(tag for (word, tag) in brown_news_tagged)
# print tag_fd.most_common()
# word_tag_pairs = nltk.bigrams(brown_news_tagged)
# noun_preceders = [a[1] for (a, b) in word_tag_pairs if b[1] == 'NOUN']
# fdist = nltk.FreqDist(noun_preceders)
# print [tag for (tag, _) in fdist.most_common()]


#################################### PART 1 ####################################

from py2neo import Node, Relationship, Graph, Path

# Phil = Node("Family",name="Phil")
# Liz = Node("Family",name="Liz")
# Chaz = Node("Family",name="Chaz")
# Anne = Node("Family",name="Anne")
# Andy = Node("Family",name="Andy")
# Ed = Node("Family",name="Ed")

# g = Graph(password="password")
# tx = g.begin()

# a = Node("Person", name="Alice")
# b = Node("Person", name="Bob")
# ab = Relationship(a, "KNOWS", b)
# tx.create(a)
# tx.create(b)
# tx.create(ab)
# tx.commit()

# print ab

# c = Node("Person", name="Carol")
# class WorksWith(Relationship): pass
# ac = WorksWith(a, c)
# print ac.type()

# s = ab | ac
# print s
# print s.nodes()
# print s.relationships()

# w = ab + Relationship(b, "LIKES", c) + ac
# print w

#################################### PART 2 ####################################

# graph = Graph(password="password")

# df01 = pd.DataFrame(graph.data("MATCH (p:Person) RETURN p.name LIMIT 10"))
# df02 = pd.DataFrame(graph.data('MATCH (p:Person {name: "Tom Hanks"}) RETURN p.name,p.born'))
# df03 = pd.DataFrame(graph.data('MATCH (m:Movie {title: "Cloud Atlas"}) RETURN m.title,m.released,m.tagline'))
# df03 = pd.DataFrame(graph.data('MATCH (m:Movie) WHERE m.title = "Cloud Atlas" RETURN m.title,m.released,m.tagline'))
# df04 = pd.DataFrame(graph.data('MATCH (m:Movie) WHERE m.released > 1990 AND m.released < 2000 RETURN m.title'))

# df05 = pd.DataFrame(graph.data('MATCH (tom:Person {name: "Tom Hanks"})-[:ACTED_IN]->(tomHanksMovies) RETURN tom,tomHanksMovies'))
# df06 = pd.DataFrame(graph.data('MATCH (cloudAtlas {title: "Cloud Atlas"})<-[:DIRECTED]-(directors) RETURN directors.name'))
# df07 = pd.DataFrame(graph.data('MATCH (tom:Person {name:"Tom Hanks"})-[:ACTED_IN]->(m)<-[:ACTED_IN]-(coActors) RETURN coActors.name'))
# df08 = pd.DataFrame(graph.data('MATCH (people:Person)-[relatedTo]-(:Movie {title: "Cloud Atlas"}) RETURN people.name, Type(relatedTo), relatedTo'))

# df09 = pd.DataFrame(graph.data('MATCH (bacon:Person {name:"Kevin Bacon"})-[*1..4]-(hollywood) RETURN DISTINCT hollywood'))
# df10 = pd.DataFrame(graph.data('MATCH p=shortestPath((bacon:Person {name:"Kevin Bacon"})-[*]-(meg:Person {name:"Meg Ryan"})) RETURN p'))
# df11 = pd.DataFrame(graph.data('MATCH (tom:Person {name:"Tom Hanks"})-[:ACTED_IN]->(m)<-[:ACTED_IN]-(coActors), (coActors)-[:ACTED_IN]->(m2)<-[:ACTED_IN]-(cocoActors) WHERE NOT (tom)-[:ACTED_IN]->(m2) RETURN cocoActors.name AS Recommended, count(*) AS Strength ORDER BY Strength DESC'))
# df12 = pd.DataFrame(graph.data('MATCH (tom:Person {name:"Tom Hanks"})-[:ACTED_IN]->(m)<-[:ACTED_IN]-(coActors), (coActors)-[:ACTED_IN]->(m2)<-[:ACTED_IN]-(cruise:Person {name:"Tom Cruise"}) RETURN tom, m, coActors, m2, cruise'))
# df13 = pd.DataFrame(graph.data('MATCH (n) RETURN n'))
# df14 = pd.DataFrame(graph.data('MATCH (a:Person),(m:Movie) OPTIONAL MATCH (a)-[r1]-(), (m)-[r2]-() DELETE a,r1,m,r2'))

#################################### PART 3 ####################################
from neo4jrestclient import client

def get_all_nodes():
    query = 'MATCH (n:Family) RETURN n'
    results = db.query(query, returns=(client.Node, str, client.Node))
    return results

def create_family_nodes(nodename,nodeborn):
    entity = db.labels.create("Family")
    for index in range(len(nodename)):
        nodeid = db.nodes.create(name=nodename[index],born=nodeborn[index])
        entity.add(nodeid)

def create_family_relationships(relationnames,relationnode1,relationnode2):
    results = get_all_nodes()
    for index in range(len(relationnames)):
        r1,r2 = '',''
        for r in results:
            if r[0]['name'] == relationnode1[index]: r1 = r[0]
            if r[0]['name'] == relationnode2[index]: r2 = r[0]
            if r1 != '' and r2 != '':
                r1.relationships.create(relationnames[index],r2)
                break    

# db = client.GraphDatabase("http://localhost:7474", username="neo4j", password="password")

# nodename = ['Phil','Liz','Chaz','Anne','Andy','Ed']
# nodeborn = [1921,1926,1948,1950,1960,1964]
# create_family_nodes(nodename,nodeborn)

# relationnode1 = ['Phil','Phil','Phil','Phil','Phil','Liz','Liz','Liz','Liz']
# relationnode2 = ['Liz','Chaz','Anne','Andy','Ed','Chaz','Anne','Andy','Ed']
# relationnames = ['MARRIED','FATHER','FATHER','FATHER','FATHER','MOTHER','MOTHER','MOTHER','MOTHER']
# create_family_relationships(relationnames,relationnode1,relationnode2)

# query = 'MATCH (f:Family)-[:MARRIED]-(s:Family) WHERE f.name="Liz" RETURN s'
# results = db.query(query, returns=(client.Node, str, client.Node))
# for r in results:
#     print r[0]['name'],r[0]['born']


#################################### TEST ####################################
import nltk
from geotext import GeoText
import parsedatetime
from datetime import datetime
from geopy.geocoders import Nominatim
from geopy.distance import great_circle

# sentence = "I want to fly from New Delhi to Mumbai next tuesday evening"

# # identify cities
# # identify location
# # identify from_loc,to_loc
# # identify airport
# # identify datetime
# # count people
# # identify return dates

# # identify cities
# def get_location_details(text):
#     places = GeoText(text)
#     data = {"cities":places.cities,"countries":places.countries}
#     return data

# # identify date
# def get_datetime_details(text):
#     p = parsedatetime.Calendar()
#     time_struct, parse_status = p.parse(text)
#     # print time_struct, parse_status
#     if parse_status != 0: return datetime(*time_struct[:6])
#     return 'error'

# # identify location
# def get_location(place):
#     geolocator = Nominatim()
#     location = geolocator.geocode(place)
#     return location.address,location.latitude,location.longitude

# # identify airport
# def get_airport(place):
#     add,lat,lon = get_location(place)
#     min_dist,closest_place = 99999,[]
#     with open('Airport.csv') as f:
#         lines = f.readlines()
#         for index,line in enumerate(lines):
#             if index == 0: continue
#             item = line.replace('\n','').split(',')
#             if item[4] != '\N':
#                 air_lat,air_lon = float(item[6]),float(item[7])
#                 distance = great_circle((lat,lon),(air_lat,air_lon))
#                 if distance < min_dist:
#                     min_dist = distance
#                     closest_place = item
#     return closest_place

# def ngrams(words,n):
#     result = []
#     for i in range(len(words)-n+1):
#         g = ' '.join(words[i:i+n])
#         result.append(g)
#     return result

# # temp = ['midnight','morning','breakfast','noon','lunch','afternoon','evening','dinner','night']
# # temp = ['today','tomorrow','friday','next monday','25th aug','11 jul','15th march next year','12/25','12.25','1/7/2020','1.7.2020']
# # temp = ['10 days later','2 weeks later','4 months later','1 year later','in 9 days','5 days from next monday','one week and one day later']

# """print get_datetime_details("this month")"""
# """print get_datetime_details("next month")"""
# """print get_datetime_details("next year")"""
# """print get_datetime_details("next week")"""

# # print get_location_details(sentence)['cities']
# # print get_datetime_details(sentence)
# # print get_airport('Delhi')

# cities = get_location_details(sentence)['cities']
# # print cities
# countries = get_location_details(sentence)['countries']
# # print countries
# for index,item in enumerate(cities): sentence = sentence.replace(cities[index],'location_'+str(index+1))
# journey = get_datetime_details(sentence)
# # print journey
# words = nltk.word_tokenize(sentence)
# # print words
# # words_tag = nltk.pos_tag(words)
# # print words_tag

# # all_grams = list(set(ngrams(words,2)+ngrams(words,3)+ngrams(words,4)+ngrams(words,5)))
# # all_datetimes = []
# # for item in all_grams:
# #     journey = get_datetime_details(item)
# #     if journey != 'error':
# #         all_datetimes.append(journey)
# # print list(set(all_datetimes))


########################################### MARKET BASKET ####################################################################################
# data = [['B2','B4','B19','B27','B39'],['B2','B4','B5','B6','B19'],['B2','B4','B5','B6','B9'],['B5','B6','B33','B35','B37'],['B5','B9','B33','B35','B37','B50'],['B5','B19','B31','B33','B35'],['B5','B6','B9','B33','B35','B39','B50'],['B5','B9','B33','B39','B42','B50'],['B5','B9','B30','B31','B35','B39'],['B5','B6','B27','B33','B39'],['B5','B9','B27','B31','B33','B37','B42'],['B5','B6','B27','B33','B35','B37','B50'],['B19','B21','B27','B31','B42'],['B19','B20','B23','B27','B42'],['B19','B23','B31','B33','B39'],['B19','B21','B27','B33','B35'],['B19','B31','B33','B39','B50'],['B20','B21','B23','B26','B27','B31','B50'],['B20','B27','B31','B37','B50'],['B20','B27','B33','B35','B37','B50'],['B23','B27','B33','B37','B39','B50'],['B31','B33','B37','B42','B50']]
# all_items_temp = []
# for item in data: all_items_temp += item
# all_items = list(set(all_items_temp))
# support_threshold = 0.1
# freuent_items = []

# def support(items):
#     valid_transactions = 0.0
#     for transaction in data:
#         count = len(items)
#         for test_item in items:
#             if test_item in transaction: count -= 1
#         if count == 0: valid_transactions += 1
#     return valid_transactions/len(data)

# def confidence(items):
#     suppa = support(items)
#     suppb = support(items[:-1])
#     return suppa/suppb

# for item in all_items:
#     if support([item])>support_threshold: frequent_items.append(item)
